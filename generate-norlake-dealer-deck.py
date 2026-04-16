"""Generate Norlake Dealer Deck — Benefits of In-House Refrigeration & Product Line Overview."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === BRAND CONFIGURATION (Norlake) ===
NAVY = RGBColor(0x00, 0x28, 0x57)
DARK_NAVY = RGBColor(0x00, 0x15, 0x32)
ACCENT_BLUE = RGBColor(0x2B, 0x7C, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF1, 0xED)
BODY_GRAY = RGBColor(0x2C, 0x3E, 0x50)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
LIGHT_BLUE_BG = RGBColor(0xEB, 0xF5, 0xFF)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
ICE_BLUE = RGBColor(0xD6, 0xEE, 0xFF)
MUTED_BLUE = RGBColor(0x8A, 0xA8, 0xC8)
SUBTLE_BLUE = RGBColor(0x5A, 0x78, 0x98)
MID_BLUE = RGBColor(0x1A, 0x4D, 0x7C)

# Norlake fonts
FONT_DISPLAY = 'Teko'
FONT_BODY = 'Trade Gothic Next'

# Margins
MARGIN_LEFT = Inches(0.75)
MARGIN_RIGHT = Inches(0.75)
CONTENT_WIDTH = Inches(11.833)

# Asset paths
ASSETS_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_WHITE = os.path.join(ASSETS_DIR, "assets", "logos", "norlake-no-oval-white.png")
LOGO_DARK = os.path.join(ASSETS_DIR, "assets", "logos", "norlake-no-oval-dark.png")

# Product images
IMG_KOLD_LOCKER = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "kold-locker-capsule-pak-left.jpg")
IMG_FAST_TRAK = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "fast-trak.jpg")
IMG_FAST_TRAK_COMBO = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "fast-trak-combo.jpg")
IMG_FINELINE = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "fineline.jpg")
IMG_CAPSULE_PAK = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "refrigeration", "capsule-pak-eco-transparent.png")
IMG_CAPSULE_PAK_LEFT = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "refrigeration", "capsule-pak-eco-left.jpg")
IMG_CONTROLLER = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "refrigeration", "controller.jpg")
IMG_LOGITEMP_LAPTOP = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "controllers", "logitemp-laptop.jpg")
IMG_LOGITEMP_TOUCH = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "controllers", "logitemp-touch.jpg")
IMG_KITCHEN = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "kitchen-setting-left.jpg")
IMG_INTERIOR = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "interior-shelving.jpg")
IMG_FOODSERVICE = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "foodservice.jpg")
IMG_INDOOR_SS = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "indoor-stainless.jpg")
IMG_ISO_9001 = os.path.join(ASSETS_DIR, "assets", "logos", "iso-9001.png")
IMG_ISO_14001 = os.path.join(ASSETS_DIR, "assets", "logos", "iso-14001.png")


# === HELPER FUNCTIONS ===

def add_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, border=None):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    return shape


def add_text(slide, text, x, y, w=None, h=None, font=FONT_BODY, size=Pt(14),
             color=BODY_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a single-line text box."""
    w = w or Inches(11)
    h = h or Inches(0.5)
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def add_multiline(slide, lines, x, y, w, h=None, font=FONT_BODY, size=Pt(12),
                  color=BODY_GRAY, bold=False, line_spacing=1.2, bullet=False):
    """Add a multi-paragraph text box from a list of strings."""
    h = h or Inches(3)
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(size.pt * (line_spacing - 1) + 2)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
    return txBox


def add_header(slide, text, y=Inches(0.4)):
    """Add a BLUF headline at the top of a content slide."""
    add_text(slide, text, x=MARGIN_LEFT, y=y, w=Inches(10), h=Inches(0.7),
             font=FONT_DISPLAY, size=Pt(40), color=NAVY, bold=True)
    add_accent_bar(slide, x=MARGIN_LEFT, y=y + Inches(0.65))


def add_accent_bar(slide, x, y, width=Inches(1.5)):
    """Blue accent rule for visual rhythm."""
    add_rect(slide, x=x, y=y, w=width, h=Pt(3), fill=ACCENT_BLUE)


def add_footer(slide, text="CONFIDENTIAL"):
    """Standard footer bar with text and logo."""
    add_rect(slide, x=Inches(0), y=Inches(6.8), w=Inches(13.333), h=Inches(0.7), fill=NAVY)
    add_text(slide, text, x=MARGIN_LEFT, y=Inches(6.95), w=Inches(5), h=Inches(0.3),
             size=Pt(8), color=WHITE)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(10.5), Inches(6.85), width=Inches(2))


def add_stat_block(slide, number, label, x, y, number_size=Pt(56)):
    """Large number + label underneath."""
    add_text(slide, number, x=x, y=y, w=Inches(2.8), h=Inches(0.8),
             font=FONT_DISPLAY, size=number_size, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, label, x=x, y=y + Inches(0.75), w=Inches(2.8), h=Inches(0.5),
             font=FONT_BODY, size=Pt(11), color=BODY_GRAY, alignment=PP_ALIGN.CENTER)


def add_card(slide, title, body_lines, x, y, w, h, accent_color=ACCENT_BLUE):
    """Card with colored top border."""
    add_rect(slide, x=x, y=y, w=w, h=Inches(0.06), fill=accent_color)
    add_rect(slide, x=x, y=y + Inches(0.06), w=w, h=h - Inches(0.06), fill=WHITE, border=LIGHT_GRAY)
    add_text(slide, title, x=x + Inches(0.2), y=y + Inches(0.15), w=w - Inches(0.4), h=Inches(0.4),
             font=FONT_DISPLAY, size=Pt(22), color=NAVY, bold=True)
    add_multiline(slide, body_lines, x=x + Inches(0.2), y=y + Inches(0.55),
                  w=w - Inches(0.4), h=h - Inches(0.65), font=FONT_BODY, size=Pt(11), color=BODY_GRAY)


# === CREATE PRESENTATION ===
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ================================================================
# SLIDE 1: TITLE — Norlake: In-House Refrigeration for Every Walk-In
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=NAVY)

if os.path.exists(LOGO_WHITE):
    s.shapes.add_picture(LOGO_WHITE, Inches(0.75), Inches(0.8), Inches(3.5))

add_text(s, 'DEALER PARTNER PRESENTATION', x=Inches(0.75), y=Inches(2.6),
         w=Inches(5), h=Inches(0.35), size=Pt(12), color=ACCENT_BLUE, bold=True)

add_text(s, 'Walk-Ins and Refrigeration\nFrom a Single Manufacturer',
         x=Inches(0.75), y=Inches(3.0), w=Inches(7), h=Inches(1.6),
         font=FONT_DISPLAY, size=Pt(56), color=WHITE, bold=True)

add_accent_bar(s, x=Inches(0.75), y=Inches(4.7), width=Inches(2))

add_text(s, 'One source for walk-ins, self-contained and remote refrigeration,\nelectronic controls, and nationwide service.',
         x=Inches(0.75), y=Inches(5.0), w=Inches(7), h=Inches(0.8),
         size=Pt(16), color=MUTED_BLUE)

# Hero product image on right
if os.path.exists(IMG_KOLD_LOCKER):
    s.shapes.add_picture(IMG_KOLD_LOCKER, Inches(7.8), Inches(1.5), width=Inches(5))

add_text(s, 'CONFIDENTIAL', x=Inches(0.75), y=Inches(6.95),
         w=Inches(5), h=Inches(0.3), size=Pt(8), color=SUBTLE_BLUE)


# ================================================================
# SLIDE 2: SECTION — Why In-House Refrigeration Matters
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=DARK_NAVY)

add_text(s, 'THE ADVANTAGE', x=Inches(0.75), y=Inches(2.2),
         w=Inches(11), h=Inches(0.4), size=Pt(13), color=ACCENT_BLUE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(s, 'Why In-House Refrigeration Matters',
         x=Inches(0.75), y=Inches(2.7), w=Inches(11.8), h=Inches(0.9),
         font=FONT_DISPLAY, size=Pt(52), color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_bar(s, x=Inches(5.6), y=Inches(3.7), width=Inches(2))

add_text(s, 'Your customers get better outcomes when the walk-in\nand the refrigeration system are engineered together.',
         x=Inches(2), y=Inches(4.2), w=Inches(9), h=Inches(0.8),
         size=Pt(18), color=MUTED_BLUE, alignment=PP_ALIGN.CENTER)

if os.path.exists(LOGO_WHITE):
    s.shapes.add_picture(LOGO_WHITE, Inches(5.4), Inches(6.2), width=Inches(2.5))


# ================================================================
# SLIDE 3: 660K+ sq ft — walk-ins and refrigeration from one source
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, '660K+ sq ft of manufacturing: walk-ins and refrigeration from one source')
add_accent_bar(s, x=MARGIN_LEFT, y=Inches(1.05))

# Left content column
bullets = [
    "Norlake manufactures both the walk-in enclosure and the refrigeration system in-house",
    "Two ISO-certified plants: Hudson, WI (ISO 9001 + 14001) and New Albany, MS (ISO 9001)",
    "Walk-in panels, self-contained units, and remote condensing units all built under one roof",
    "One engineering team designs the enclosure and the refrigeration together for optimized performance",
]
add_multiline(s, bullets, x=MARGIN_LEFT, y=Inches(1.4),
              w=Inches(6.5), h=Inches(4.5), size=Pt(13), color=BODY_GRAY,
              line_spacing=1.6)

# Right side — kitchen lifestyle image
if os.path.exists(IMG_KITCHEN):
    s.shapes.add_picture(IMG_KITCHEN, Inches(7.8), Inches(1.4), width=Inches(5))

# ISO badges at bottom
if os.path.exists(IMG_ISO_9001):
    s.shapes.add_picture(IMG_ISO_9001, Inches(0.75), Inches(5.6), width=Inches(0.9))
if os.path.exists(IMG_ISO_14001):
    s.shapes.add_picture(IMG_ISO_14001, Inches(1.8), Inches(5.6), width=Inches(0.9))

add_footer(s)


# ================================================================
# SLIDE 4: Single-source refrigeration eliminates finger-pointing
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Single-source refrigeration eliminates finger-pointing and speeds service')
add_accent_bar(s, x=MARGIN_LEFT, y=Inches(1.05))

# Two-column comparison
col_w = Inches(5.5)
col_gap = Inches(0.5)
col1_x = MARGIN_LEFT
col2_x = MARGIN_LEFT + col_w + col_gap

# Left column: Typical Setup
add_rect(s, col1_x, Inches(1.5), col_w, Inches(0.45), fill=RGBColor(0xFE, 0xF2, 0xF2))
add_text(s, 'TYPICAL: SEPARATE VENDORS', x=col1_x + Inches(0.2), y=Inches(1.52),
         w=col_w, h=Inches(0.4), size=Pt(13), color=RED, bold=True)

typical_bullets = [
    "Walk-in from one manufacturer, refrigeration from another",
    "Warranty disputes when a problem spans both systems",
    "Two vendors to coordinate during installation",
    "Finger-pointing delays service calls by days or weeks",
    "No guarantee the systems are optimized for each other",
]
add_multiline(s, typical_bullets, x=col1_x + Inches(0.2), y=Inches(2.1),
              w=col_w - Inches(0.4), h=Inches(3.5), size=Pt(12), color=BODY_GRAY,
              line_spacing=1.5)

# Right column: Norlake
add_rect(s, col2_x, Inches(1.5), col_w, Inches(0.45), fill=GREEN_BG)
add_text(s, 'NORLAKE: ONE MANUFACTURER', x=col2_x + Inches(0.2), y=Inches(1.52),
         w=col_w, h=Inches(0.4), size=Pt(13), color=GREEN, bold=True)

norlake_bullets = [
    "Walk-in and refrigeration engineered and built together",
    "One warranty, one call, one resolution",
    "Single point of contact for quoting, ordering, and service",
    "Factory-matched systems ship ready to install",
    "Nationwide Approved Service Provider network backs every unit",
]
add_multiline(s, norlake_bullets, x=col2_x + Inches(0.2), y=Inches(2.1),
              w=col_w - Inches(0.4), h=Inches(3.5), size=Pt(12), color=BODY_GRAY,
              line_spacing=1.5)

add_footer(s)


# ================================================================
# SLIDE 5: STATS — Key Numbers
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT_GRAY)
add_header(s, 'The numbers your customers care about')
add_accent_bar(s, x=MARGIN_LEFT, y=Inches(1.05))

stat_y = Inches(1.8)
stat_gap = Inches(3.0)
start_x = Inches(0.5)

add_stat_block(s, '14,000+', 'stocked walk-in configurations', x=start_x, y=stat_y)
add_stat_block(s, '2-DAY', 'shipping on stock orders', x=start_x + stat_gap, y=stat_y)
add_stat_block(s, '~50%', 'energy savings with Capsule Pak ECO\nvs. legacy R-404A (verified)', x=start_x + stat_gap * 2, y=stat_y)
add_stat_block(s, '15-YEAR', 'structural panel warranty', x=start_x + stat_gap * 3, y=stat_y)

# Second row
stat_y2 = Inches(3.8)
add_stat_block(s, '660K+', 'sq ft of manufacturing\nacross 2 plants', x=start_x, y=stat_y2)
add_stat_block(s, 'R-290', 'natural refrigerant (GWP of 3)\n3 years ahead of AIM Act', x=start_x + stat_gap, y=stat_y2)
add_stat_block(s, '18-MONTH', 'parts and labor warranty\n(50% longer than industry standard)', x=start_x + stat_gap * 2, y=stat_y2)
add_stat_block(s, 'EST. 1947', 'trusted by restaurants, grocers,\nand institutions nationwide', x=start_x + stat_gap * 3, y=stat_y2)

add_footer(s)


# ================================================================
# SLIDE 6: SECTION — The Product Line
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, DARK_NAVY)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=DARK_NAVY)

add_text(s, 'PRODUCT LINE OVERVIEW', x=Inches(0.75), y=Inches(2.2),
         w=Inches(11), h=Inches(0.4), size=Pt(13), color=ACCENT_BLUE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(s, 'Walk-Ins, Refrigeration, and Controls',
         x=Inches(0.75), y=Inches(2.7), w=Inches(11.8), h=Inches(0.9),
         font=FONT_DISPLAY, size=Pt(52), color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_bar(s, x=Inches(5.6), y=Inches(3.7), width=Inches(2))

add_text(s, 'From stock walk-ins to custom configurations, self-contained\nto remote refrigeration, standard to smart controls.',
         x=Inches(2), y=Inches(4.2), w=Inches(9), h=Inches(0.8),
         size=Pt(18), color=MUTED_BLUE, alignment=PP_ALIGN.CENTER)

if os.path.exists(LOGO_WHITE):
    s.shapes.add_picture(LOGO_WHITE, Inches(5.4), Inches(6.2), width=Inches(2.5))


# ================================================================
# SLIDE 7: Kold Locker
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Kold Locker: 14,000+ configurations ship in 2 business days')
add_accent_bar(s, x=MARGIN_LEFT, y=Inches(1.05))

kl_bullets = [
    "Broadest standard walk-in range in the industry: 14,000+ stocked configurations",
    "2-day shipping eliminates project delays for your customers",
    "Foamed-in-place polyurethane insulation with foam-to-foam joints (no thermal bridging)",
    "Cam-lock assembly for fast, reliable field installation",
    "15-year structural panel warranty; 18-month parts and labor",
    "Indoor and outdoor models for restaurants, grocery, institutional, floral, beer caves",
]
add_multiline(s, kl_bullets, x=MARGIN_LEFT, y=Inches(1.4),
              w=Inches(6), h=Inches(4.5), size=Pt(13), color=BODY_GRAY,
              line_spacing=1.5)

# Hero image right
if os.path.exists(IMG_KOLD_LOCKER):
    s.shapes.add_picture(IMG_KOLD_LOCKER, Inches(7.5), Inches(1.2), width=Inches(5.3))

add_footer(s)


# ================================================================
# SLIDE 8: Capsule Pak ECO
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Capsule Pak ECO: R-290 self-contained refrigeration cuts energy ~50%')
add_accent_bar(s, x=MARGIN_LEFT, y=Inches(1.05))

cp_bullets = [
    "Patented self-contained system (U.S. Patent No. 11,859,885)",
    "R-290 natural refrigerant: GWP of 3 vs. 3,922 for R-404A (99.9% reduction)",
    "Verified energy savings: 56% cooler, 51% freezer vs. legacy systems",
    "Pre-charged, cord-and-plug design: most models 115V, no field brazing",
    "Up to 28% shorter and 20% narrower than legacy units",
    "LogiTemp electronic controller standard on every unit",
]
add_multiline(s, cp_bullets, x=MARGIN_LEFT, y=Inches(1.4),
              w=Inches(6.5), h=Inches(4.5), size=Pt(13), color=BODY_GRAY,
              line_spacing=1.5)

# Energy comparison table area
tbl_y = Inches(4.6)
add_text(s, 'VERIFIED ENERGY PERFORMANCE', x=MARGIN_LEFT, y=tbl_y,
         w=Inches(4), h=Inches(0.35), size=Pt(10), color=NAVY, bold=True)

# Table rows
row_data = [
    ('Cooler (kWh/day)', '5.6', '2.5', '56%'),
    ('Freezer (kWh/day)', '11.3', '5.6', '51%'),
]
headers = ['', 'Legacy R-404A', 'Capsule Pak ECO', 'Savings']
col_x = [MARGIN_LEFT, MARGIN_LEFT + Inches(1.5), MARGIN_LEFT + Inches(3.0), MARGIN_LEFT + Inches(4.5)]
col_w = [Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.0)]

# Header row
for j, hdr in enumerate(headers):
    add_text(s, hdr, x=col_x[j], y=tbl_y + Inches(0.35), w=col_w[j], h=Inches(0.3),
             size=Pt(9), color=NAVY, bold=True)

for i, (label, old, new, sav) in enumerate(row_data):
    ry = tbl_y + Inches(0.65) + Inches(i * 0.3)
    vals = [label, old, new, sav]
    for j, val in enumerate(vals):
        clr = GREEN if j == 3 else BODY_GRAY
        add_text(s, val, x=col_x[j], y=ry, w=col_w[j], h=Inches(0.25),
                 size=Pt(9), color=clr, bold=(j == 3))

# Hero image right — transparent Capsule Pak
if os.path.exists(IMG_CAPSULE_PAK):
    s.shapes.add_picture(IMG_CAPSULE_PAK, Inches(8), Inches(1.2), width=Inches(4.5))

add_footer(s)


# ================================================================
# SLIDE 9: Split-Pak Remote Systems
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Split-Pak: remote systems from 1/2 HP to 15 HP cover every application')
add_accent_bar(s, x=MARGIN_LEFT, y=Inches(1.05))

sp_bullets = [
    "Remote condensing units for small reach-through coolers to large cold storage",
    "Pre-charged quick-connect line sets up to 15 HP / 50 ft reduce field brazing",
    "Pre-wired, factory-mounted components with traceable tubing connections",
    "EC motors standard on single-phase models",
    "All 2026+ quotes use AIM Act-compliant refrigerants (R-454A, R-454C)",
]
add_multiline(s, sp_bullets, x=MARGIN_LEFT, y=Inches(1.4),
              w=Inches(6.5), h=Inches(3), size=Pt(13), color=BODY_GRAY,
              line_spacing=1.5)

# Split-Pak ECO callout box
add_rect(s, MARGIN_LEFT, Inches(4.2), Inches(6.5), Inches(1.8), fill=LIGHT_BLUE_BG)
add_text(s, 'SPLIT-PAK ECO', x=MARGIN_LEFT + Inches(0.2), y=Inches(4.3),
         w=Inches(4), h=Inches(0.35), size=Pt(12), color=NAVY, bold=True)
add_multiline(s, [
    "First R-290 remote refrigeration system in the industry",
    "Extends natural refrigerant benefits beyond the 1.75 HP self-contained limit",
    "Your customers get R-290 energy savings on larger walk-in applications",
], x=MARGIN_LEFT + Inches(0.2), y=Inches(4.7),
    w=Inches(6), h=Inches(1.2), size=Pt(11), color=BODY_GRAY, line_spacing=1.4)

# Image right side
if os.path.exists(IMG_FOODSERVICE):
    s.shapes.add_picture(IMG_FOODSERVICE, Inches(7.8), Inches(1.2), width=Inches(5))

add_footer(s)


# ================================================================
# SLIDE 10: LogiTemp Electronic Controllers
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'LogiTemp: smart controls save up to 27% energy and log 30 days of data')
add_accent_bar(s, x=MARGIN_LEFT, y=Inches(1.05))

lt_bullets = [
    "Patented electronic controller: standard on all Capsule Pak ECO units",
    "Demand Defrost defrosts only when needed (not on a fixed timer)",
    "30-day room and coil temperature data logging built in",
    "Wireless or Cat5 remote monitoring and programming via Sitrad",
    "Up to 27% energy savings vs. all-mechanical control systems",
]
add_multiline(s, lt_bullets, x=MARGIN_LEFT, y=Inches(1.4),
              w=Inches(6.5), h=Inches(3), size=Pt(13), color=BODY_GRAY,
              line_spacing=1.5)

# Reverse Cycle Defrost callout
add_rect(s, MARGIN_LEFT, Inches(4.2), Inches(6.5), Inches(1.8), fill=GREEN_BG)
add_text(s, 'LOGITEMP PLUS: REVERSE CYCLE DEFROST', x=MARGIN_LEFT + Inches(0.2), y=Inches(4.3),
         w=Inches(5), h=Inches(0.35), size=Pt(12), color=GREEN, bold=True)
add_multiline(s, [
    "Up to 80% less defrost energy vs. electric heater defrost",
    "Freezer defrost: 3-5 minutes vs. 20-30 minutes (85% faster)",
    "Available on Split-Pak 6 HP+ systems",
], x=MARGIN_LEFT + Inches(0.2), y=Inches(4.7),
    w=Inches(6), h=Inches(1.2), size=Pt(11), color=BODY_GRAY, line_spacing=1.4)

# Image right side
if os.path.exists(IMG_LOGITEMP_LAPTOP):
    s.shapes.add_picture(IMG_LOGITEMP_LAPTOP, Inches(7.8), Inches(1.2), width=Inches(5))

add_footer(s)


# ================================================================
# SLIDE 11: Fast-Trak and FineLine
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Fast-Trak and FineLine extend the line from semi-custom to premium')
add_accent_bar(s, x=MARGIN_LEFT, y=Inches(1.05))

# Two cards side by side
card_w = Inches(5.5)
card_h = Inches(4.2)
card_gap = Inches(0.6)
card_y = Inches(1.5)
card1_x = MARGIN_LEFT
card2_x = MARGIN_LEFT + card_w + card_gap

# Fast-Trak card
add_card(s, 'FAST-TRAK',
         [
             "Semi-custom and custom walk-ins for projects beyond Kold Locker standard options",
             "Non-standard dimensions, specialized layouts, and complex configurations",
             "Same foamed-in-place polyurethane insulation and build quality as Kold Locker",
             "Fills the gap between stock and fully custom",
         ],
         x=card1_x, y=card_y, w=card_w, h=card_h, accent_color=ACCENT_BLUE)

# FineLine card
add_card(s, 'FINELINE',
         [
             "Premium architectural walk-in line for design-sensitive applications",
             "Aesthetics and design integration alongside full performance",
             "Ideal for front-of-house, visible installations, and upscale environments",
             "Same engineering, same warranty, same in-house refrigeration options",
         ],
         x=card2_x, y=card_y, w=card_w, h=card_h, accent_color=GOLD)

# Product images below cards
if os.path.exists(IMG_FAST_TRAK):
    s.shapes.add_picture(IMG_FAST_TRAK, card1_x + Inches(0.5), card_y + Inches(3.0),
                         width=Inches(2.2))

if os.path.exists(IMG_FINELINE):
    s.shapes.add_picture(IMG_FINELINE, card2_x + Inches(0.5), card_y + Inches(3.0),
                         width=Inches(2.2))

add_footer(s)


# ================================================================
# SLIDE 12: CLOSER — Partner with Norlake
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=NAVY)

if os.path.exists(LOGO_WHITE):
    s.shapes.add_picture(LOGO_WHITE, Inches(4.7), Inches(0.8), Inches(4))

add_text(s, 'One Call. One Warranty.\nOne Manufacturer.',
         x=Inches(1.5), y=Inches(2.6), w=Inches(10), h=Inches(1.5),
         font=FONT_DISPLAY, size=Pt(56), color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_accent_bar(s, x=Inches(5.6), y=Inches(4.2), width=Inches(2))

add_multiline(s, [
    "Walk-ins, self-contained and remote refrigeration, and electronic controls",
    "all engineered and built in-house.",
    "",
    "Contact your Norlake sales representative for pricing and availability.",
    "800-955-5253  |  norlake.com",
], x=Inches(2), y=Inches(4.5), w=Inches(9), h=Inches(2),
    size=Pt(16), color=MUTED_BLUE, line_spacing=1.3)

add_text(s, 'CONFIDENTIAL', x=Inches(0.75), y=Inches(6.95),
         w=Inches(5), h=Inches(0.3), size=Pt(8), color=SUBTLE_BLUE)


# === SAVE ===
output_path = r"c:\Users\andre\Desktop\Norlake-Dealer-Deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
