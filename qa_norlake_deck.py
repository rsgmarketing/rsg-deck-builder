"""Visual QA: bounds check + wireframe render for Norlake Dealer Deck."""

from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image, ImageDraw, ImageFont
import os

PPTX_PATH = r"C:\Users\andre\Desktop\Norlake-Dealer-Deck.pptx"
OUTPUT_DIR = r"C:\tmp\rsg-deck-builder\qa_wireframes"

os.makedirs(OUTPUT_DIR, exist_ok=True)


def qa_deck(pptx_path):
    """Visual QA: bounds check + wireframe render for every slide."""
    prs = Presentation(pptx_path)
    issues = []
    wireframes = []

    # Scale: 1333x750 canvas represents 13.333" x 7.5"
    SCALE = 100  # pixels per inch
    W, H = 1333, 750
    SAFETY_Y = 650  # 6.5" in pixels

    for slide_idx, slide in enumerate(prs.slides):
        # Create wireframe canvas
        img = Image.new('RGB', (W, H), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        shapes_data = []
        for shape in slide.shapes:
            x = int(shape.left / 914400 * SCALE)
            y = int(shape.top / 914400 * SCALE)
            w = int(shape.width / 914400 * SCALE)
            h = int(shape.height / 914400 * SCALE)
            shapes_data.append((x, y, w, h, shape))

            # Draw shape fill
            fill_color = (220, 220, 220)  # default gray
            try:
                if shape.fill and shape.fill.fore_color:
                    rgb = shape.fill.fore_color.rgb
                    fill_color = (rgb[0], rgb[1], rgb[2])
            except:
                pass

            if hasattr(shape, 'image'):
                try:
                    _ = shape.image
                    # Image placeholder - draw with blue border
                    draw.rectangle([x, y, x + w, y + h], outline=(43, 124, 204), width=2)
                    draw.text((x + 4, y + 4), '[IMG]', fill=(43, 124, 204))
                except:
                    draw.rectangle([x, y, x + w, y + h], fill=fill_color, outline=(180, 180, 180))
            else:
                draw.rectangle([x, y, x + w, y + h], fill=fill_color, outline=(180, 180, 180))

            # Draw text content
            if shape.has_text_frame:
                text = shape.text_frame.text[:100]
                if text.strip():
                    # Estimate font size
                    font_size = 10
                    try:
                        for p in shape.text_frame.paragraphs:
                            if p.runs:
                                font_size = int(p.runs[0].font.size / 12700) if p.runs[0].font.size else 10
                                break
                    except:
                        pass

                    text_color = (44, 62, 80)
                    try:
                        for p in shape.text_frame.paragraphs:
                            if p.runs and p.runs[0].font.color and p.runs[0].font.color.rgb:
                                rgb = p.runs[0].font.color.rgb
                                text_color = (rgb[0], rgb[1], rgb[2])
                                break
                    except:
                        pass

                    # Draw text (approximate)
                    try:
                        fnt = ImageFont.truetype("segoeui.ttf", max(8, min(font_size, 36)))
                    except:
                        fnt = ImageFont.load_default()

                    # Wrap text to fit within the shape
                    display_text = text[:80]
                    draw.text((x + 4, y + 4), display_text, fill=text_color, font=fnt)

            # --- BOUNDS CHECKS ---
            bottom = y + h
            if bottom > SAFETY_Y and y < SAFETY_Y:
                # Content crosses into footer zone
                issues.append(f"Slide {slide_idx + 1}: Shape at y={y/SCALE:.1f}\" extends to {bottom/SCALE:.1f}\" (below 6.5\" safety line)")
                draw.rectangle([x, y, x + w, y + h], outline=(255, 0, 0), width=3)

        # Check for overlaps between non-background shapes
        for i in range(len(shapes_data)):
            for j in range(i + 1, len(shapes_data)):
                x1, y1, w1, h1, s1 = shapes_data[i]
                x2, y2, w2, h2, s2 = shapes_data[j]
                # Skip full-slide backgrounds
                if w1 > W * 0.9 or w2 > W * 0.9:
                    continue
                # Skip footer-zone shapes (y >= 6.5")
                if y1 >= SAFETY_Y or y2 >= SAFETY_Y:
                    continue
                # Check overlap
                if (x1 < x2 + w2 and x1 + w1 > x2 and y1 < y2 + h2 and y1 + h1 > y2):
                    overlap_area = (min(x1+w1, x2+w2) - max(x1, x2)) * (min(y1+h1, y2+h2) - max(y1, y2))
                    if overlap_area > 500:  # ignore tiny/intentional overlaps
                        # Check if one is intentionally inside the other (card pattern)
                        # If one shape fully contains the other, skip (it's a card with content)
                        contained = (x1 <= x2 and y1 <= y2 and x1+w1 >= x2+w2 and y1+h1 >= y2+h2) or \
                                    (x2 <= x1 and y2 <= y1 and x2+w2 >= x1+w1 and y2+h2 >= y1+h1)
                        if not contained:
                            # Get text snippets for identification
                            t1 = s1.text_frame.text[:30] if s1.has_text_frame else "[shape]"
                            t2 = s2.text_frame.text[:30] if s2.has_text_frame else "[shape]"
                            issues.append(f"Slide {slide_idx + 1}: Overlap between \"{t1}\" at ({x1/SCALE:.1f}\",{y1/SCALE:.1f}\") and \"{t2}\" at ({x2/SCALE:.1f}\",{y2/SCALE:.1f}\")")

        # Draw safety line
        for sx in range(0, W, 10):
            draw.line([(sx, SAFETY_Y), (min(sx + 5, W), SAFETY_Y)], fill=(255, 0, 0), width=2)
        draw.text((W - 130, SAFETY_Y - 15), "6.5\" SAFETY", fill=(255, 0, 0))

        # Draw slide number
        draw.text((10, H - 20), f"Slide {slide_idx + 1}", fill=(100, 100, 100))

        wireframe_path = os.path.join(OUTPUT_DIR, f"slide_{slide_idx + 1}_wireframe.png")
        img.save(wireframe_path)
        wireframes.append(wireframe_path)

    return issues, wireframes


# Run QA
issues, wireframes = qa_deck(PPTX_PATH)

if issues:
    print("=== ISSUES FOUND ===")
    for issue in issues:
        print(f"  - {issue}")
else:
    print("=== NO ISSUES FOUND ===")

print(f"\nWireframes saved ({len(wireframes)} slides):")
for wf in wireframes:
    print(f"  {wf}")
print("\nReview each wireframe image before delivering the deck.")
