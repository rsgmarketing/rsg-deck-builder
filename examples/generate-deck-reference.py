"""Generate Fast-Trak Strategy v3 PowerPoint — Two Paths to ~2-Week Ship."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === COLOR PALETTE (from Norlake Consultant Intro reference deck) ===
NAVY = RGBColor(0x00, 0x28, 0x57)
DARK_NAVY = RGBColor(0x0B, 0x13, 0x20)
MID_BLUE = RGBColor(0x1A, 0x4D, 0x7C)
BODY_GRAY = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF1, 0xED)
ACCENT_BLUE = RGBColor(0x2B, 0x7C, 0xCC)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1A, 0x7A, 0x4A)
ORANGE = RGBColor(0xD4, 0x78, 0x0A)
ICE_BLUE = RGBColor(0xD6, 0xEE, 0xFF)
MUTED_BLUE = RGBColor(0x8A, 0xA8, 0xC8)
DARK_PANEL = RGBColor(0x1D, 0x2D, 0x42)
SUBTLE_BLUE = RGBColor(0x5A, 0x78, 0x98)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF2, 0xFC)
GREEN_BG = RGBColor(0xED, 0xF7, 0xF1)
RED_BG = RGBColor(0xFD, 0xF0, 0xEE)
ORANGE_BG = RGBColor(0xFE, 0xF6, 0xEC)
CC_BLUE = RGBColor(0xCC, 0xDD, 0xEE)
PURPLE = RGBColor(0x6C, 0x3D, 0x91)
PURPLE_BG = RGBColor(0xF3, 0xEC, 0xF8)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
TEAL_BG = RGBColor(0xE6, 0xF5, 0xF5)

FONT_DISPLAY = 'Teko'
FONT_BODY = 'Trade Gothic Next'

logo_white = r'C:\Users\andre\Desktop\Claude Code\Fast Trak\extracted_images\slide1_Picture 3.png'
logo_dark = r'C:\Users\andre\Desktop\Claude Code\Fast Trak\extracted_images\slide2_Picture 6.png'

# Owner badge colors
OWNER_COLORS = {
    'MANUFACTURING': (RGBColor(0x8B, 0x45, 0x13), RGBColor(0xFD, 0xF0, 0xE2)),
    'SALES': (GREEN, GREEN_BG),
    'FINANCE': (PURPLE, PURPLE_BG),
    'PRODUCT': (ACCENT_BLUE, LIGHT_BLUE_BG),
    'OPERATIONS': (TEAL, TEAL_BG),
    'ECOMMERCE': (ORANGE, ORANGE_BG),
    'ENGINEERING': (MID_BLUE, LIGHT_BLUE_BG),
    'MARKETING': (RGBColor(0xB0, 0x30, 0x60), RGBColor(0xFD, 0xEC, 0xF2)),
    'COMPETITIVE INTEL': (BODY_GRAY, LIGHT_GRAY),
    'STRATEGIC DECISION': (NAVY, RGBColor(0xE0, 0xE8, 0xF0)),
    'KEY ACCOUNT': (RGBColor(0xD4, 0x78, 0x0A), ORANGE_BG),
    'SUPPLY CHAIN': (TEAL, TEAL_BG),
    'DATA PULL': (ACCENT_BLUE, LIGHT_BLUE_BG),
}


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def add_text(slide, left, top, width, height, text,
             font_name=FONT_BODY, size=Pt(14), color=BODY_GRAY,
             bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_name=FONT_BODY,
                  size=Pt(12), color=BODY_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                  line_spacing=Pt(18)):
    """Add text with explicit line breaks as separate paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        if line_spacing:
            p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
    return txBox


def add_header(slide, text, top=Inches(0.6), left=Inches(0.8)):
    return add_text(slide, left, top, Inches(11), Inches(0.8), text,
                    font_name=FONT_DISPLAY, size=Pt(40), color=NAVY, bold=True)


def add_accent_bar(slide, top=Inches(1.2), left=Inches(0.8)):
    return add_rect(slide, left, top, Inches(1.5), Pt(3), fill_color=ACCENT_BLUE)


def add_subheader(slide, text, top=Inches(1.35), left=Inches(0.8)):
    return add_text(slide, left, top, Inches(11), Inches(0.4), text,
                    font_name=FONT_BODY, size=Pt(15), color=MID_BLUE, bold=True)


def add_logo_white(slide):
    slide.shapes.add_picture(logo_white, Inches(9.5), Inches(6.5), Inches(3.0))


def add_logo_dark(slide):
    slide.shapes.add_picture(logo_dark, Inches(10.0), Inches(6.55), Inches(2.5))


def add_footer(slide, dark=False):
    c = WHITE if dark else BODY_GRAY
    add_text(slide, Inches(0.8), Inches(7.0), Inches(5), Inches(0.3),
             'INTERNAL \u2014 CONFIDENTIAL  |  February 2026',
             font_name=FONT_BODY, size=Pt(9), color=c)


def add_owner_badge(slide, x, y, owner_text):
    """Add a small colored badge indicating question owner."""
    key = owner_text.upper()
    # Try to match to known colors, fallback to default
    fg_color = BODY_GRAY
    bg_color = LIGHT_GRAY
    for k, (fg, bg) in OWNER_COLORS.items():
        if k in key:
            fg_color = fg
            bg_color = bg
            break
    badge_w = Inches(0.12 * len(owner_text) + 0.15)
    if badge_w > Inches(2.0):
        badge_w = Inches(2.0)
    add_rect(slide, x, y, badge_w, Inches(0.2), fill_color=bg_color)
    add_text(slide, x + Inches(0.05), y - Inches(0.01), badge_w - Inches(0.1), Inches(0.2),
             owner_text.upper(), size=Pt(7), color=fg_color, bold=True)
    return badge_w


def add_question_row(slide, x, y, q_num, q_text, owner, width=Inches(11.0)):
    """Add a formatted question row with number bubble, text, and owner badge."""
    # Number bubble
    add_rect(slide, x, y + Inches(0.02), Inches(0.3), Inches(0.3), fill_color=ACCENT_BLUE)
    num_str = str(q_num)
    add_text(slide, x + Inches(0.01), y - Inches(0.01), Inches(0.3), Inches(0.32),
             num_str, font_name=FONT_DISPLAY, size=Pt(16), color=WHITE,
             bold=True, alignment=PP_ALIGN.CENTER)
    # Question text
    add_text(slide, x + Inches(0.4), y, width - Inches(2.5), Inches(0.35),
             q_text, size=Pt(11), color=BODY_GRAY)
    # Owner badge
    badge_x = x + width - Inches(1.8)
    add_owner_badge(slide, badge_x, y + Inches(0.05), owner)


def add_question_block(slide, x, y, q_num, q_text, owner, width=Inches(5.0)):
    """Add a question with number, text, and owner badge in a card layout."""
    # Number
    add_text(slide, x, y, Inches(0.4), Inches(0.25),
             f'Q{q_num}', font_name=FONT_DISPLAY, size=Pt(16), color=ACCENT_BLUE, bold=True)
    # Owner badge on same line
    add_owner_badge(slide, x + width - Inches(1.5), y + Inches(0.02), owner)
    # Question text below
    add_text(slide, x, y + Inches(0.28), width, Inches(0.5),
             q_text, size=Pt(11), color=BODY_GRAY)


# ================================================================
# CREATE PRESENTATION
# ================================================================
prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
blank = prs.slide_layouts[6]

# ================================================================
# SLIDE 1: TITLE (dark)
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill_color=NAVY)
s.shapes.add_picture(logo_white, Inches(1.0), Inches(1.0), Inches(4.0))
add_text(s, Inches(1.0), Inches(2.6), Inches(4), Inches(0.35),
         'INTERNAL PLANNING FRAMEWORK', size=Pt(11), color=ACCENT_BLUE, bold=True)
add_text(s, Inches(1.0), Inches(3.0), Inches(10), Inches(1.2),
         'Fast-Trak Product Line Strategy',
         font_name=FONT_DISPLAY, size=Pt(56), color=WHITE, bold=True)
add_text(s, Inches(1.0), Inches(4.3), Inches(10), Inches(0.6),
         'An Internal Planning Framework',
         size=Pt(22), color=ICE_BLUE)
add_rect(s, Inches(1.0), Inches(5.0), Inches(2.0), Pt(2), fill_color=ACCENT_BLUE)
add_text(s, Inches(1.0), Inches(5.3), Inches(10), Inches(0.8),
         'This document asks the questions we need to answer\nbefore committing to a direction.',
         size=Pt(14), color=MUTED_BLUE)
add_text(s, Inches(1.0), Inches(6.8), Inches(5), Inches(0.3),
         'CONFIDENTIAL  |  February 2026  |  v3', size=Pt(9), color=SUBTLE_BLUE)

# ================================================================
# SLIDE 2: HOW TO USE THIS DOCUMENT
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'How to Use This Document')
add_accent_bar(s)

add_text(s, Inches(0.8), Inches(1.7), Inches(11), Inches(0.8),
         'This deck is structured around questions, not conclusions. Each section presents a hypothesis and then asks the questions needed to validate or invalidate it. Every question has an owner and a response type.',
         size=Pt(14), color=BODY_GRAY)

# Response type legend
types_y = Inches(2.8)
add_text(s, Inches(0.8), types_y, Inches(3), Inches(0.3),
         'RESPONSE TYPES', size=Pt(11), color=NAVY, bold=True)

type_items = [
    ('DATA PULL', 'Requires pulling order data, financial records, or system reports.', ACCENT_BLUE, LIGHT_BLUE_BG),
    ('YES / NO', 'A direct operational or feasibility question.', GREEN, GREEN_BG),
    ('SALES INPUT', 'Requires qualitative input from customer-facing teams.', RGBColor(0x1A, 0x7A, 0x4A), GREEN_BG),
    ('DECISION', 'Requires a strategic call from leadership.', NAVY, RGBColor(0xE0, 0xE8, 0xF0)),
]
for i, (label, desc, fg, bg) in enumerate(type_items):
    y = types_y + Inches(0.4) + Inches(i * 0.55)
    add_rect(s, Inches(0.8), y, Inches(1.3), Inches(0.25), fill_color=bg)
    add_text(s, Inches(0.85), y - Inches(0.02), Inches(1.2), Inches(0.25),
             label, size=Pt(9), color=fg, bold=True)
    add_text(s, Inches(2.3), y, Inches(4), Inches(0.25),
             desc, size=Pt(10), color=BODY_GRAY)

# Owner legend
owners_y = Inches(2.8)
add_text(s, Inches(6.5), owners_y, Inches(3), Inches(0.3),
         'QUESTION OWNERS', size=Pt(11), color=NAVY, bold=True)

owner_items = ['Manufacturing', 'Sales', 'Finance', 'Product', 'Operations', 'Ecommerce', 'Marketing']
for i, owner in enumerate(owner_items):
    y = owners_y + Inches(0.4) + Inches(i * 0.42)
    add_owner_badge(s, Inches(6.5), y, owner)

# Goal
add_rect(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8), fill_color=LIGHT_BLUE_BG)
add_text(s, Inches(1.0), Inches(5.85), Inches(11), Inches(0.7),
         'GOAL: Answer 43 questions across 8 domains. Two paths to ~2-week ship. Build consensus through data and direct answers, not top-down recommendations.',
         size=Pt(13), color=NAVY, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 3: THE PROBLEM — LEAD TIME IDENTITY CRISIS
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'The Current State: An Identity Crisis')
add_accent_bar(s)

by = Inches(2.0)
bl = Inches(2.5)

# KL bar
add_text(s, Inches(0.8), by, Inches(1.6), Inches(0.45),
         'KOLD LOCKER', size=Pt(12), color=GREEN, bold=True, alignment=PP_ALIGN.RIGHT)
add_rect(s, bl, by + Inches(0.05), Inches(0.5), Inches(0.35), fill_color=GREEN)
add_text(s, bl + Inches(0.6), by, Inches(2), Inches(0.45),
         '2 days', font_name=FONT_DISPLAY, size=Pt(20), color=GREEN, bold=True)

# FT bar
fy = by + Inches(0.7)
add_text(s, Inches(0.8), fy, Inches(1.6), Inches(0.45),
         'FAST-TRAK', size=Pt(12), color=RED, bold=True, alignment=PP_ALIGN.RIGHT)
add_rect(s, bl, fy + Inches(0.05), Inches(8.0), Inches(0.35), fill_color=RED)
add_text(s, bl + Inches(8.2), fy, Inches(2), Inches(0.45),
         '4\u201310 weeks', font_name=FONT_DISPLAY, size=Pt(20), color=RED, bold=True)

# Fineline bar
fly = fy + Inches(0.7)
add_text(s, Inches(0.8), fly, Inches(1.6), Inches(0.45),
         'FINELINE', size=Pt(12), color=RGBColor(0xE0, 0x50, 0x40), bold=True, alignment=PP_ALIGN.RIGHT)
add_rect(s, bl, fly + Inches(0.05), Inches(8.0), Inches(0.35), fill_color=RGBColor(0xE0, 0x50, 0x40))
add_text(s, bl + Inches(8.2), fly, Inches(2), Inches(0.45),
         '4\u201310 weeks', font_name=FONT_DISPLAY, size=Pt(20), color=RGBColor(0xE0, 0x50, 0x40), bold=True)

add_text(s, Inches(7.5), fly + Inches(0.5), Inches(3), Inches(0.3),
         '\u2191 IDENTICAL LEAD TIMES', size=Pt(12), color=RED, bold=True, alignment=PP_ALIGN.RIGHT)

# Callouts
cy = Inches(4.6)
callouts = [
    ('~2%', 'The price gap between Fast-Trak and\nFineline. Confusing, not invisible.', Inches(0.8)),
    ('\u201cFast\u201d Is in the Name', 'But the product isn\u2019t fast. It ships at\nthe same factory lead time as custom.', Inches(4.8)),
    ('1-ft Increments', 'Factory stocks 2-ft panels. Not leverageable\nfor stock speed. Process speed possible?', Inches(8.8)),
]
for title, desc, x in callouts:
    add_rect(s, x, cy, Inches(3.5), Inches(1.3), fill_color=RED_BG)
    fs = Pt(36) if title == '~2%' else Pt(22)
    ty = cy + Inches(0.1) if title == '~2%' else cy + Inches(0.15)
    add_text(s, x + Inches(0.2), ty, Inches(3.1), Inches(0.5),
             title, font_name=FONT_DISPLAY, size=fs, color=RED, bold=True)
    add_text(s, x + Inches(0.2), cy + Inches(0.65), Inches(3.1), Inches(0.5),
             desc, size=Pt(11), color=BODY_GRAY)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 4: THE PROBLEM — PRICE PERCEPTION (NEW)
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'The Price Reality')
add_accent_bar(s)
add_subheader(s, 'Price sensitivity varies dramatically by job type and channel')

# Two-column layout
# Left: New Construction
add_rect(s, Inches(0.8), Inches(2.0), Inches(5.3), Inches(2.2), fill_color=GREEN_BG, line_color=GREEN)
add_text(s, Inches(1.0), Inches(2.05), Inches(4.9), Inches(0.3),
         'NEW CONSTRUCTION', size=Pt(12), color=GREEN, bold=True)
add_text(s, Inches(1.0), Inches(2.4), Inches(4.9), Inches(1.6),
         'Walk-in is a fraction of total project cost. Schedule and lead time often matter more than a 5-10% delta. Architect/GC drives the spec. Buyers have more flexibility on price.',
         size=Pt(12), color=BODY_GRAY)

# Right: Replacement / Retrofit
add_rect(s, Inches(6.5), Inches(2.0), Inches(5.8), Inches(2.2), fill_color=RED_BG, line_color=RED)
add_text(s, Inches(6.7), Inches(2.05), Inches(5.4), Inches(0.3),
         'REPLACEMENT / RETROFIT', size=Pt(12), color=RED, bold=True)
add_text(s, Inches(6.7), Inches(2.4), Inches(5.4), Inches(1.6),
         'Operator is paying directly. Every dollar matters. Downtime = lost revenue. Speed AND price both critical. These buyers shop aggressively and compare line-by-line.',
         size=Pt(12), color=BODY_GRAY)

# Ecommerce callout
add_rect(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.0), fill_color=ORANGE_BG)
add_text(s, Inches(1.0), Inches(4.55), Inches(2), Inches(0.3),
         'ECOMMERCE', size=Pt(12), color=ORANGE, bold=True)
add_text(s, Inches(1.0), Inches(4.85), Inches(11), Inches(0.6),
         'Ecommerce buyers key hard on pricing. Walk-ins are perceived as commodities in this channel, even if that\u2019s not our positioning. Price visibility is total \u2014 buyers compare across brands in seconds.',
         size=Pt(12), color=BODY_GRAY)

# Bottom insight
add_rect(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8), fill_color=NAVY)
add_text(s, Inches(1.0), Inches(5.85), Inches(11), Inches(0.7),
         'The ~2% FT-to-Fineline gap isn\u2019t just \u201cinvisible\u201d \u2014 it\u2019s actively confusing. If the products are priced the same and ship the same, what is the buyer paying for? This is a credibility problem, not just a marketing problem.',
         size=Pt(12), color=WHITE)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 5: QUESTIONS — DO WE HAVE A REAL PROBLEM?
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Questions: Do We Have a Real Problem?')
add_accent_bar(s)
add_subheader(s, 'Before fixing Fast-Trak, confirm the diagnosis')

questions_5 = [
    (1, 'What % of FT orders could have been built as Fineline with no configuration change?', 'Data Pull'),
    (2, 'How often do dealers ask "what\'s the difference between FT and Fineline?"', 'Sales'),
    (3, 'Have we lost deals because FT lead time matched Fineline \u2014 i.e., buyer went to a faster competitor?', 'Sales'),
    (4, 'What is the current volume split (units + revenue) across KL / FT / Fineline?', 'Finance'),
]
for i, (num, q, owner) in enumerate(questions_5):
    y = Inches(2.0) + Inches(i * 0.95)
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.8), fill_color=LIGHT_BLUE_BG if i % 2 == 0 else WHITE)
    add_question_row(s, Inches(0.9), y + Inches(0.2), num, q, owner, width=Inches(11.3))

add_text(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.5),
         'If the answer to Q1 is >50% and Q4 shows FT is a small fraction of revenue, the case for change is strong. If FT is a large revenue line with differentiated demand, the calculus shifts.',
         size=Pt(11), color=MID_BLUE, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 6: TWO PATHS TO ~2-WEEK SHIP (dark)
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)

add_text(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
         'Two Paths to ~2-Week Ship', font_name=FONT_DISPLAY, size=Pt(40), color=WHITE, bold=True)
add_rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Pt(3), fill_color=ACCENT_BLUE)

# Path A card (green)
add_rect(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.6), fill_color=DARK_PANEL, line_color=GREEN)
add_text(s, Inches(1.0), Inches(1.65), Inches(2), Inches(0.25),
         'PATH A', size=Pt(10), color=GREEN, bold=True)
add_text(s, Inches(1.0), Inches(1.95), Inches(5.1), Inches(0.5),
         'Speed Through Inventory', font_name=FONT_DISPLAY, size=Pt(28), color=GREEN, bold=True)
add_text(s, Inches(1.0), Inches(2.5), Inches(5.1), Inches(0.6),
         'Same 2-ft stock panels as KL. Assembly-only workflow. Panels on shelf \u2014 just add options and assemble.',
         size=Pt(12), color=CC_BLUE)
add_text(s, Inches(1.0), Inches(3.15), Inches(5.1), Inches(0.25),
         'TRADE-OFF: Loses 1-ft sizing. FT grid = KL grid.', size=Pt(10), color=MUTED_BLUE)
add_text(s, Inches(1.0), Inches(3.45), Inches(5.1), Inches(0.25),
         'SPEED FROM: Panels on shelf + assembly scheduling.', size=Pt(10), color=GREEN)

# Path B card (blue)
add_rect(s, Inches(6.7), Inches(1.6), Inches(5.6), Inches(2.6), fill_color=DARK_PANEL, line_color=ACCENT_BLUE)
add_text(s, Inches(6.9), Inches(1.65), Inches(2), Inches(0.25),
         'PATH B', size=Pt(10), color=ACCENT_BLUE, bold=True)
add_text(s, Inches(6.9), Inches(1.95), Inches(5.2), Inches(0.5),
         'Speed Through Process', font_name=FONT_DISPLAY, size=Pt(28), color=ACCENT_BLUE, bold=True)
add_text(s, Inches(6.9), Inches(2.5), Inches(5.2), Inches(0.6),
         'Every FT config = pre-defined SKU with pre-engineered BOM. Zero engineering at order time. Panels built to order.',
         size=Pt(12), color=CC_BLUE)
add_text(s, Inches(6.9), Inches(3.15), Inches(5.2), Inches(0.25),
         'TRADE-OFF: Panels built to order. Requires Ops 2-week commitment.', size=Pt(10), color=MUTED_BLUE)
add_text(s, Inches(6.9), Inches(3.45), Inches(5.2), Inches(0.25),
         'SPEED FROM: No engineering/quoting + production optimization.', size=Pt(10), color=ACCENT_BLUE)

# Proposed lead time bars
add_text(s, Inches(0.8), Inches(4.5), Inches(5), Inches(0.3),
         'PROPOSED LEAD TIME STRUCTURE', size=Pt(11), color=ACCENT_BLUE, bold=True)

bars = [
    ('KOLD LOCKER', Inches(0.5), GREEN, '2 days'),
    ('FAST-TRAK', Inches(2.5), ACCENT_BLUE, '~2 weeks'),
    ('FINELINE', Inches(8.0), SUBTLE_BLUE, '4\u201310 weeks'),
]
for i, (label, bw, color, desc) in enumerate(bars):
    y = Inches(4.95) + Inches(i * 0.5)
    add_text(s, Inches(0.8), y, Inches(1.6), Inches(0.4),
             label, size=Pt(10), color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    add_rect(s, Inches(2.6), y + Inches(0.05), bw, Inches(0.3), fill_color=color)
    add_text(s, Inches(2.6) + bw + Inches(0.15), y, Inches(3), Inches(0.4),
             desc, size=Pt(12), color=color, bold=True)

# Callout
add_rect(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5), fill_color=DARK_PANEL)
add_text(s, Inches(1.0), Inches(6.42), Inches(11), Inches(0.45),
         'Pre-engineering is a prerequisite for both paths. The manufacturing model is proven \u2014 the question is which mechanism gets us back there.',
         size=Pt(12), color=CC_BLUE, bold=True)

add_logo_white(s)
add_footer(s, dark=True)

# ================================================================
# SLIDE 7: QUESTIONS — PATH A: MANUFACTURING FEASIBILITY
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Path A: Manufacturing Feasibility')
add_accent_bar(s)
add_subheader(s, 'Can we build from stock 2-ft panels in 2 weeks?')

mfg_qs = [
    (5, 'Can we assemble a FT unit from stock 2-ft panels + staged components in 10 business days?', 'Manufacturing'),
    (6, 'Could filler panels (small bridging panels) make odd-foot sizes buildable from stock without stocking full 1-ft panels?', 'Engineering'),
    (7, 'Do we have the plant capacity to build and stock filler panels if that approach is viable?', 'Manufacturing'),
    (8, 'What is the assembly time delta between a standard KL unit and a combo/remote-refrig unit?', 'Manufacturing'),
    (9, 'Why was the original rapid-build model abandoned? What specifically would need to change to make it work again?', 'Operations'),
]
for i, (num, q, owner) in enumerate(mfg_qs):
    y = Inches(2.0) + Inches(i * 0.88)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.75), fill_color=bg)
    add_question_row(s, Inches(0.9), y + Inches(0.18), num, q, owner, width=Inches(11.3))

add_rect(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4), fill_color=ORANGE_BG)
add_text(s, Inches(1.0), Inches(6.32), Inches(11), Inches(0.35),
         'Q9 is the #1 feasibility input. If we don\u2019t understand why it was abandoned, we risk repeating the same failure.',
         size=Pt(11), color=ORANGE, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 8: QUESTIONS — PATH A: SIZING & INVENTORY
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Path A: Sizing & Inventory')
add_accent_bar(s)
add_subheader(s, 'Can we serve FT demand from stock 2-ft panels?')

sizing_qs = [
    (10, 'What % of current FT orders are already in 2-ft increments (and could be built from stock panels)?', 'Data Pull'),
    (11, 'What % of FT orders use odd-foot sizes (5\u2019, 7\u2019, 9\u2019, 11\u2019)?', 'Data Pull'),
    (12, 'For odd-foot FT orders, what price premium do they carry over the nearest 2-ft equivalent?', 'Data Pull'),
    (13, 'What additional component inventory (refrigeration systems, combo hardware, female bottom rails) would need to be staged?', 'Supply Chain'),
    (14, 'What is the carrying cost of that incremental inventory?', 'Finance'),
]
for i, (num, q, owner) in enumerate(sizing_qs):
    y = Inches(2.0) + Inches(i * 0.88)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.75), fill_color=bg)
    add_question_row(s, Inches(0.9), y + Inches(0.18), num, q, owner, width=Inches(11.3))

add_rect(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4), fill_color=GREEN_BG)
add_text(s, Inches(1.0), Inches(6.32), Inches(11), Inches(0.35),
         'If Q10 shows >80% of FT orders are already 2-ft, the case for a stock-panel FT program is very strong.',
         size=Pt(11), color=GREEN, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 9: QUESTIONS — PATH B: SKU & BOM FEASIBILITY
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Path B: SKU & BOM Feasibility')
add_accent_bar(s)
add_subheader(s, 'Can we pre-engineer every FT configuration as a known SKU?')

pathb_sku_qs = [
    (34, 'How many distinct valid FT configurations exist today? What is the total SKU count?', 'Engineering / Product'),
    (35, 'Can we pre-engineer BOMs for all valid FT configs? What is the effort and timeline?', 'Engineering'),
    (36, 'What is the current average engineering/quoting time per FT order?', 'Operations'),
    (37, 'Can FT configs be reduced to a manageable set (<500 SKUs) without losing significant demand?', 'Product / Engineering'),
]
for i, (num, q, owner) in enumerate(pathb_sku_qs):
    y = Inches(2.0) + Inches(i * 0.95)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.8), fill_color=bg)
    add_question_row(s, Inches(0.9), y + Inches(0.2), num, q, owner, width=Inches(11.3))

add_rect(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.6), fill_color=LIGHT_BLUE_BG)
add_text(s, Inches(1.0), Inches(5.85), Inches(11), Inches(0.5),
         'If Q34 shows a manageable config count and Q35 confirms BOM pre-engineering is feasible, Path B is a real option. If SKU count is unmanageable, Path A may be the only viable route.',
         size=Pt(11), color=MID_BLUE, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 10: QUESTIONS — PATH B: OPERATIONS PROCESS
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Path B: Operations Process')
add_accent_bar(s)
add_subheader(s, 'Can Ops deliver a 2-week build for pre-configured SKUs?')

pathb_ops_qs = [
    (38, 'With a pre-built BOM and zero engineering, what is the realistic production lead time today?', 'Operations / Mfg'),
    (39, 'What process steps add time beyond panel production + assembly? Which can be shortened for pre-configured SKUs?', 'Operations'),
    (40, 'Can dedicated production slots / priority queuing for pre-configured FT hit 2 weeks reliably? Volume threshold?', 'Operations / Mfg'),
    (41, 'What components could be pre-staged even if panels are built to order? What is the staging cost?', 'Supply Chain / Finance'),
]
for i, (num, q, owner) in enumerate(pathb_ops_qs):
    y = Inches(2.0) + Inches(i * 0.95)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.8), fill_color=bg)
    add_question_row(s, Inches(0.9), y + Inches(0.2), num, q, owner, width=Inches(11.3))

add_rect(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.6), fill_color=ORANGE_BG)
add_text(s, Inches(1.0), Inches(5.85), Inches(11), Inches(0.5),
         'Q38 is the anchor: if current production time (minus engineering) is already close to 2 weeks, Path B is viable with process optimization.',
         size=Pt(11), color=ORANGE, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 11: QUESTIONS — PATH A vs. PATH B DECISION
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Path A vs. Path B Decision')
add_accent_bar(s)
add_subheader(s, 'Once both paths are assessed, the strategic question: which one (or both)?')

compare_qs = [
    (42, 'If both paths are feasible, which has lower total operational cost? (Inventory carrying for A vs. process investment for B.)', 'Finance / Operations'),
    (43, 'Could a hybrid work \u2014 stock 2-ft (Path A) for high-volume, pre-configured SKUs (Path B) for 1-ft/odd \u2014 to get both speed and flexibility?', 'Strategic Decision'),
]
for i, (num, q, owner) in enumerate(compare_qs):
    y = Inches(2.0) + Inches(i * 1.2)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(1.0), fill_color=bg)
    add_question_row(s, Inches(0.9), y + Inches(0.3), num, q, owner, width=Inches(11.3))

# Summary box
add_rect(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(1.2), fill_color=NAVY)
add_text(s, Inches(1.0), Inches(4.65), Inches(11), Inches(1.1),
         'Path A is simpler but loses 1-ft sizing. Path B retains sizing but requires deeper operational commitment. A hybrid could offer the best of both \u2014 if the complexity is manageable.\n\nThis is the central strategic question after both paths are validated.',
         size=Pt(13), color=WHITE, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 12: QUESTIONS — FINELINE BOUNDARY
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Questions: What Must Stay in Fineline?')
add_accent_bar(s)

# Proposed boundary box
add_text(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.3),
         'PROPOSED FINELINE BOUNDARY', size=Pt(11), color=NAVY, bold=True)
add_rect(s, Inches(0.8), Inches(2.05), Inches(5.5), Inches(2.4), fill_color=LIGHT_GRAY, line_color=NAVY)
fineline_items = [
    '\u2022 Non-standard dimensions (not 2-ft increments) *Path A only',
    '\u2022 Glass doors or custom door styles',
    '\u2022 Custom finishes (beyond standard stucco)',
    '\u2022 Heavy-duty / structural floors (>800 lb/sq ft)',
    '\u2022 Water-cooled refrigeration',
    '\u2022 3+ compartments',
    '\u2022 Non-rectangular shapes or angles',
    '\u2022 5\u201d or 6\u201d panel thickness',
]
add_multiline(s, Inches(1.0), Inches(2.1), Inches(5.1), Inches(2.3),
              fineline_items, size=Pt(10), color=BODY_GRAY)

# Questions on right
qs = [
    (15, 'Which options/features should ALWAYS push to Fineline regardless of speed goals?', 'Product'),
    (16, 'Are there current FT options not on this list that cannot be rapid-shipped?', 'Manufacturing'),
    (17, 'Are there Fineline options that COULD move to a rapid-ship model?', 'Manufacturing'),
]
for i, (num, q, owner) in enumerate(qs):
    y = Inches(1.9) + Inches(i * 1.0)
    add_rect(s, Inches(6.8), y, Inches(5.5), Inches(0.85), fill_color=LIGHT_BLUE_BG if i % 2 == 0 else WHITE)
    add_question_block(s, Inches(7.0), y + Inches(0.1), num, q, owner, width=Inches(5.1))

add_text(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
         'The goal: a crisp, defensible boundary. Everything above the line = Fineline. Everything below = potentially rapid-shippable.',
         size=Pt(12), color=MID_BLUE, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 10: QUESTIONS — KL-TO-FT BOUNDARY
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Questions: KL-to-FT Boundary')
add_accent_bar(s)
add_subheader(s, 'What triggers should push a Kold Locker order to Fast-Trak?')

# Current triggers
add_text(s, Inches(0.8), Inches(2.0), Inches(5.3), Inches(0.3),
         'CURRENT KL \u2192 FT TRIGGERS', size=Pt(11), color=NAVY, bold=True)
triggers = [
    ('Remote Refrigeration', 'CP Remote or Split-Pak', 'PRIMARY', GREEN),
    ('Combo / Dual Compartment', 'Cooler-freezer combinations', 'KEY', ACCENT_BLUE),
    ('Female Bottom Rail', 'Required for certain installs', 'KEY', ACCENT_BLUE),
    ('Additional Heights', '6\u20197\u201d, 8\u20197\u201d, 8\u20194\u201d floorless', 'SECONDARY', MID_BLUE),
    ('Odd-Foot Sizes', '5\u2019, 7\u2019, 9\u2019, 11\u2019 widths/depths', 'UNQUANTIFIED', ORANGE),
]
for i, (name, desc, badge, color) in enumerate(triggers):
    y = Inches(2.35) + Inches(i * 0.55)
    add_rect(s, Inches(0.8), y, Inches(0.9), Inches(0.22), fill_color=color)
    add_text(s, Inches(0.83), y - Inches(0.02), Inches(0.85), Inches(0.22),
             badge, size=Pt(7), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(1.8), y - Inches(0.02), Inches(2), Inches(0.22),
             name, size=Pt(10), color=NAVY, bold=True)
    add_text(s, Inches(3.8), y - Inches(0.02), Inches(2.2), Inches(0.22),
             desc, size=Pt(9), color=BODY_GRAY)

# Questions on right
boundary_qs = [
    (18, 'Are remote refrigeration, combos, and female bottom rail the right FT triggers? Are there others?', 'Product'),
    (19, 'Should 8\u20197\u201d height go into KL (if demand warrants) or stay as an FT driver?', 'Sales'),
    (20, 'What % of KL-to-FT escalations are ONLY because of odd-foot sizing vs. actual option needs?', 'Data Pull'),
]
for i, (num, q, owner) in enumerate(boundary_qs):
    y = Inches(2.0) + Inches(i * 1.15)
    add_rect(s, Inches(6.5), y, Inches(5.8), Inches(1.0), fill_color=LIGHT_BLUE_BG if i % 2 == 0 else WHITE)
    add_question_block(s, Inches(6.7), y + Inches(0.12), num, q, owner, width=Inches(5.4))

add_rect(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.6), fill_color=NAVY)
add_text(s, Inches(1.0), Inches(5.55), Inches(11), Inches(0.5),
         'Critical: Remote refrigeration and combos must stay OFF Kold Locker to preserve the FT tier boundary. This is a business decision, not a technical one.',
         size=Pt(11), color=WHITE, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 11: QUESTIONS — PRICING STRATEGY
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Questions: Pricing Strategy')
add_accent_bar(s)
add_subheader(s, 'Price matters more than we\u2019ve been assuming')

# Context box
add_rect(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0), fill_color=RED_BG)
add_text(s, Inches(1.0), Inches(2.05), Inches(11), Inches(0.9),
         'Correction from v1: Price is NOT the least important factor. For replacement/retrofit jobs and ecommerce, price is often the #1 or #2 decision driver. Buyers treat walk-ins as commodities. The ~2% FT-to-Fineline gap is actively confusing because it suggests sameness.',
         size=Pt(12), color=RED)

pricing_qs = [
    (21, 'For replacement/retrofit jobs, how price-sensitive are buyers compared to new construction?', 'Sales'),
    (22, 'What is the ecommerce price elasticity \u2014 do we lose meaningful volume at a 5% premium? 10%?', 'Ecommerce'),
    (23, 'Should FT command a speed premium over Fineline for comparable configs? (Expedited = premium in most industries.)', 'Strategic Decision'),
    (24, 'How do competitors price their quick-ship vs. custom tiers?', 'Competitive Intel'),
]
for i, (num, q, owner) in enumerate(pricing_qs):
    y = Inches(3.3) + Inches(i * 0.82)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.7), fill_color=bg)
    add_question_row(s, Inches(0.9), y + Inches(0.16), num, q, owner, width=Inches(11.3))

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 12: QUESTIONS — CHANNEL & ECOMMERCE
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Questions: Channel & Ecommerce')
add_accent_bar(s)
add_subheader(s, 'Where the rubber meets the road')

channel_qs = [
    (25, 'What is WebstaurantStore\u2019s appetite for \u201cShips in 2 Weeks\u201d Fast-Trak SKUs?', 'Key Account'),
    (26, 'What FT configurations currently sell online, and at what premium over KL?', 'Data Pull'),
    (27, 'Would a guaranteed ship date increase FT ecommerce conversion rate?', 'Ecommerce'),
    (28, 'How do large traditional dealers (non-ecommerce) view the FT vs. Fineline distinction today?', 'Sales'),
]
for i, (num, q, owner) in enumerate(channel_qs):
    y = Inches(2.0) + Inches(i * 1.0)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.85), fill_color=bg)
    add_question_row(s, Inches(0.9), y + Inches(0.22), num, q, owner, width=Inches(11.3))

# Insight box
add_rect(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.7), fill_color=GREEN_BG)
add_text(s, Inches(1.0), Inches(6.05), Inches(11), Inches(0.6),
         'Remote refrigeration \u2014 the primary KL\u2192FT driver \u2014 currently has no quick-ship path on ecommerce. A 2-week FT with remote refrig could be a significant ecommerce unlock.',
         size=Pt(12), color=GREEN, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 13: THE NAMING QUESTION (dark)
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)

add_text(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
         'Should We Rename Fast-Trak?', font_name=FONT_DISPLAY, size=Pt(40), color=WHITE, bold=True)
add_rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Pt(3), fill_color=ACCENT_BLUE)

add_text(s, Inches(0.8), Inches(1.7), Inches(11), Inches(0.8),
         'The name \u201cFast-Trak\u201d only works if the product is actually fast. That\u2019s either a problem to fix or a name to change.',
         size=Pt(16), color=CC_BLUE)

# Two scenario boxes
# Scenario 1: Speed restored
add_rect(s, Inches(0.8), Inches(2.8), Inches(5.3), Inches(3.2), fill_color=DARK_PANEL, line_color=GREEN)
add_text(s, Inches(1.0), Inches(2.85), Inches(4.9), Inches(0.3),
         'IF SPEED IS RESTORED (~2 WEEK SHIP)', size=Pt(11), color=GREEN, bold=True)
add_text(s, Inches(1.0), Inches(3.2), Inches(4.9), Inches(0.4),
         'The name becomes an asset.',
         size=Pt(16), color=WHITE, bold=True)
add_text(s, Inches(1.0), Inches(3.65), Inches(4.9), Inches(2.0),
         '\u201cFast-Trak\u201d means fast when it genuinely ships in 2 weeks vs. 6\u20138 for Fineline. Brand equity is restored. But we should still consider whether a name better aligned to the \u201cKold Locker Plus\u201d positioning might outperform.',
         size=Pt(12), color=MUTED_BLUE)

# Scenario 2: Speed not restored
add_rect(s, Inches(6.5), Inches(2.8), Inches(5.8), Inches(3.2), fill_color=DARK_PANEL, line_color=RED)
add_text(s, Inches(6.7), Inches(2.85), Inches(5.4), Inches(0.3),
         'IF SPEED IS NOT RESTORED', size=Pt(11), color=RED, bold=True)
add_text(s, Inches(6.7), Inches(3.2), Inches(5.4), Inches(0.4),
         'The name is a liability.',
         size=Pt(16), color=WHITE, bold=True)
add_text(s, Inches(6.7), Inches(3.65), Inches(5.4), Inches(2.0),
         '\u201cFast-Trak\u201d at factory lead time is a broken promise. Every time a dealer quotes it and the lead time is 6+ weeks, the name undermines our credibility. Renaming becomes essential, not optional.',
         size=Pt(12), color=MUTED_BLUE)

add_text(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4),
         'The following slides present naming options for both scenarios.',
         size=Pt(13), color=ACCENT_BLUE)

add_logo_white(s)
add_footer(s, dark=True)

# ================================================================
# SLIDE 14: NAMING OPTIONS — SPEED RESTORED
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Naming Options: If Speed Is Restored')
add_accent_bar(s)
add_subheader(s, 'Three paths for a product that actually ships fast')

# Option A: Keep Fast-Trak
ax = Inches(0.8)
aw = Inches(3.6)
ay = Inches(2.0)
add_rect(s, ax, ay, aw, Inches(4.2), line_color=GREEN)
add_rect(s, ax, ay, aw, Inches(0.4), fill_color=GREEN)
add_text(s, ax + Inches(0.1), ay + Inches(0.05), aw - Inches(0.2), Inches(0.3),
         'OPTION A: STATUS QUO', size=Pt(10), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, ax + Inches(0.2), ay + Inches(0.55), aw - Inches(0.4), Inches(0.5),
         'Keep\n\u201cFast-Trak\u201d', font_name=FONT_DISPLAY, size=Pt(28), color=NAVY, bold=True)
add_text(s, ax + Inches(0.2), ay + Inches(1.3), aw - Inches(0.4), Inches(0.2),
         'RATIONALE', size=Pt(9), color=GREEN, bold=True)
add_text(s, ax + Inches(0.2), ay + Inches(1.5), aw - Inches(0.4), Inches(1.2),
         'Name becomes honest with 2-week ship. Existing brand equity and dealer familiarity preserved. Zero transition cost. No channel confusion.',
         size=Pt(11), color=BODY_GRAY)
add_text(s, ax + Inches(0.2), ay + Inches(2.9), aw - Inches(0.4), Inches(0.2),
         'RISK', size=Pt(9), color=RED, bold=True)
add_text(s, ax + Inches(0.2), ay + Inches(3.1), aw - Inches(0.4), Inches(0.8),
         'Doesn\u2019t signal strategic change to market. \u201cFast-Trak\u201d still sounds generic.',
         size=Pt(11), color=BODY_GRAY)

# Option B: Kold Locker Pro
bx = Inches(4.7)
add_rect(s, bx, ay, aw, Inches(4.2), line_color=ACCENT_BLUE)
add_rect(s, bx, ay, aw, Inches(0.4), fill_color=ACCENT_BLUE)
add_text(s, bx + Inches(0.1), ay + Inches(0.05), aw - Inches(0.2), Inches(0.3),
         'OPTION B: LEVERAGE KL BRAND', size=Pt(10), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, bx + Inches(0.2), ay + Inches(0.55), aw - Inches(0.4), Inches(0.5),
         '\u201cKold Locker\nPro\u201d', font_name=FONT_DISPLAY, size=Pt(28), color=NAVY, bold=True)
add_text(s, bx + Inches(0.2), ay + Inches(1.3), aw - Inches(0.4), Inches(0.2),
         'RATIONALE', size=Pt(9), color=GREEN, bold=True)
add_text(s, bx + Inches(0.2), ay + Inches(1.5), aw - Inches(0.4), Inches(1.2),
         'Leverages KL\u2019s strong brand equity. Instant dealer understanding: \u201cit\u2019s a KL with more options.\u201d \u201cPro\u201d implies upgrade. Clean family hierarchy.',
         size=Pt(11), color=BODY_GRAY)
add_text(s, bx + Inches(0.2), ay + Inches(2.9), aw - Inches(0.4), Inches(0.2),
         'RISK', size=Pt(9), color=RED, bold=True)
add_text(s, bx + Inches(0.2), ay + Inches(3.1), aw - Inches(0.4), Inches(0.8),
         'Could cannibalize KL. Loses established \u201cFast-Trak\u201d recognition. May imply KL is \u201cnon-Pro.\u201d',
         size=Pt(11), color=BODY_GRAY)

# Option C: QuikBild
cx = Inches(8.6)
add_rect(s, cx, ay, aw, Inches(4.2), line_color=MID_BLUE)
add_rect(s, cx, ay, aw, Inches(0.4), fill_color=MID_BLUE)
add_text(s, cx + Inches(0.1), ay + Inches(0.05), aw - Inches(0.2), Inches(0.3),
         'OPTION C: HERITAGE PLAY', size=Pt(10), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, cx + Inches(0.2), ay + Inches(0.55), aw - Inches(0.4), Inches(0.5),
         '\u201cQuikBild\u201d', font_name=FONT_DISPLAY, size=Pt(28), color=NAVY, bold=True)
add_text(s, cx + Inches(0.2), ay + Inches(1.3), aw - Inches(0.4), Inches(0.2),
         'RATIONALE', size=Pt(9), color=GREEN, bold=True)
add_text(s, cx + Inches(0.2), ay + Inches(1.5), aw - Inches(0.4), Inches(1.2),
         'Echoes Master-Bilt\u2019s 10Bilt heritage. Communicates speed directly. Distinctive in market. Signals a return to the rapid-build model.',
         size=Pt(11), color=BODY_GRAY)
add_text(s, cx + Inches(0.2), ay + Inches(2.9), aw - Inches(0.4), Inches(0.2),
         'RISK', size=Pt(9), color=RED, bold=True)
add_text(s, cx + Inches(0.2), ay + Inches(3.1), aw - Inches(0.4), Inches(0.8),
         'New name = transition cost. May not resonate if dealers don\u2019t know 10Bilt history. Unconventional spelling.',
         size=Pt(11), color=BODY_GRAY)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 15: NAMING OPTIONS — SPEED NOT RESTORED
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Naming Options: If Speed Is Not Restored')
add_accent_bar(s)
add_subheader(s, 'Four options that don\u2019t make a promise we can\u2019t keep')

names = [
    ('D', 'FlexLine', 'Emphasizes flexibility. Rhymes with Fineline for family cohesion. No speed promise.', 'Sounds similar to Fineline \u2014 could merge in dealer minds.', ACCENT_BLUE),
    ('E', 'SelectLine', 'Implies curated options. Pairs with Fineline. Professional, clean.', 'Generic. Doesn\u2019t differentiate strongly from Fineline.', MID_BLUE),
    ('F', 'ProLine', 'Premium positioning. Works across all channels. No speed claim.', 'Very generic. ProLine exists in other industries. Weak differentiation.', TEAL),
    ('G', 'OptiKold', 'Distinctive. Uses Norlake\u2019s \u201cKold\u201d branding. Implies optimized solution.', 'Unfamiliar. Requires education. May feel forced.', PURPLE),
]
for i, (letter, name, rationale, risk, color) in enumerate(names):
    x = Inches(0.6) + Inches(i * 3.05)
    y = Inches(2.0)
    w = Inches(2.85)
    add_rect(s, x, y, w, Inches(4.2), line_color=color)
    add_rect(s, x, y, w, Inches(0.35), fill_color=color)
    add_text(s, x + Inches(0.1), y + Inches(0.02), w - Inches(0.2), Inches(0.3),
             f'OPTION {letter}', size=Pt(10), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(0.45), w - Inches(0.3), Inches(0.45),
             f'\u201c{name}\u201d', font_name=FONT_DISPLAY, size=Pt(26), color=NAVY, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(1.1), w - Inches(0.3), Inches(0.15),
             'RATIONALE', size=Pt(8), color=GREEN, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(1.25), w - Inches(0.3), Inches(1.2),
             rationale, size=Pt(10), color=BODY_GRAY)
    add_text(s, x + Inches(0.15), y + Inches(2.7), w - Inches(0.3), Inches(0.15),
             'RISK', size=Pt(8), color=RED, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(2.85), w - Inches(0.3), Inches(1.0),
             risk, size=Pt(10), color=BODY_GRAY)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 16: QUESTIONS — NAMING DECISION
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Questions: The Naming Decision')
add_accent_bar(s)
add_subheader(s, 'What do we need to know before choosing?')

naming_qs = [
    (29, 'Does the Fast-Trak name carry meaningful equity with dealers that\u2019s worth preserving?', 'Sales'),
    (30, 'If we rename, what is the full transition cost? (Collateral, systems, dealer re-education, ecommerce listings.)', 'Marketing'),
    (31, 'Would a name change signal positive strategic momentum to the market, or create confusion during transition?', 'Sales'),
]
for i, (num, q, owner) in enumerate(naming_qs):
    y = Inches(2.2) + Inches(i * 1.1)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.9), fill_color=bg)
    add_question_row(s, Inches(0.9), y + Inches(0.25), num, q, owner, width=Inches(11.3))

# Decision framework
add_rect(s, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.0), fill_color=NAVY)
add_text(s, Inches(1.0), Inches(5.65), Inches(11), Inches(0.9),
         'Decision framework: If we restore speed \u2192 keep \u201cFast-Trak\u201d (lowest risk) or consider \u201cKL Pro\u201d (highest strategic upside). If we cannot restore speed \u2192 rename is essential. \u201cFlexLine\u201d or \u201cSelectLine\u201d are safest; \u201cOptiKold\u201d is boldest.',
         size=Pt(13), color=WHITE, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 17: COMPETITIVE CONTEXT
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Competitive Context')
add_accent_bar(s)
add_subheader(s, 'What the market looks like')

# Competitor lead time bars
comp_y = Inches(2.0)
competitors = [
    ('KPS GLOBAL', '<48 hrs', Inches(0.4), GREEN),
    ('ARCTIC', '2 days', Inches(0.4), GREEN),
    ('AMERIKOOLER', '3 days', Inches(0.5), GREEN),
    ('KOLD LOCKER', '2 days', Inches(0.4), ACCENT_BLUE),
    ('FAST-TRAK (today)', '4\u201310 weeks', Inches(7.0), RED),
    ('FINELINE', '4\u201310 weeks', Inches(7.0), RGBColor(0xA0, 0xA0, 0xA0)),
]
for i, (name, time, bw, color) in enumerate(competitors):
    y = comp_y + Inches(i * 0.52)
    is_ours = name in ('KOLD LOCKER', 'FAST-TRAK (today)', 'FINELINE')
    text_color = color if is_ours else BODY_GRAY
    add_text(s, Inches(0.8), y, Inches(2.0), Inches(0.4),
             name, size=Pt(10), color=text_color, bold=is_ours, alignment=PP_ALIGN.RIGHT)
    add_rect(s, Inches(3.0), y + Inches(0.05), bw, Inches(0.3), fill_color=color)
    add_text(s, Inches(3.0) + bw + Inches(0.15), y, Inches(2), Inches(0.4),
             time, size=Pt(10), color=color, bold=True)

# Questions
comp_qs = [
    (32, 'What is Master-Bilt\u2019s current Ready-Bilt lead time? Same problem, or have they solved it differently?', 'Competitive Intel'),
    (33, 'How do KPS Global, Amerikooler, and Arctic position their quick-ship vs. custom tiers, and how do they price the gap?', 'Competitive Intel'),
]
for i, (num, q, owner) in enumerate(comp_qs):
    y = Inches(5.2) + Inches(i * 0.85)
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.7), fill_color=bg)
    add_question_row(s, Inches(0.9), y + Inches(0.15), num, q, owner, width=Inches(11.3))

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 18: IMPLEMENTATION ROADMAP
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Implementation: 18-Week Roadmap')
add_accent_bar(s)
add_subheader(s, 'Each phase answers a cluster of questions')

phases = [
    ('PHASE 1', 'Data &\nFeasibility', 'Weeks 1\u20134', GREEN,
     ['Q5\u2013Q14 + Q34\u2013Q41', '\u2022 Path A: Mfg assessment', '\u2022 Path B: SKU count + BOMs', '\u2022 Path B: Ops lead time', '\u2022 Pull FT order data']),
    ('PHASE 2', 'Boundaries &\nPath Decision', 'Weeks 4\u20138', ACCENT_BLUE,
     ['Q15\u2013Q20 + Q42\u2013Q43', '\u2022 Path A vs B decision', '\u2022 Fineline boundary', '\u2022 KL-to-FT triggers', '\u2022 Pricing structure']),
    ('PHASE 3', 'Manufacturing\n& Systems', 'Weeks 8\u201316', MID_BLUE,
     ['Execution', '\u2022 Stage inventory', '\u2022 Pre-engineer BOMs', '\u2022 Production slots', '\u2022 Automated order entry']),
    ('PHASE 4', 'Collateral &\nChannel', 'Weeks 12\u201318', NAVY,
     ['Answers Q25\u2013Q31', '\u2022 Resolve inconsistencies', '\u2022 Reposition brochures', '\u2022 Naming decision', '\u2022 Train sales + dealers']),
]
for i, (phase, title, weeks, color, items) in enumerate(phases):
    x = Inches(0.6) + Inches(i * 3.1)
    y = Inches(1.8)
    w = Inches(2.9)
    add_rect(s, x, y, w, Inches(0.35), fill_color=color)
    add_text(s, x + Inches(0.1), y + Inches(0.03), w - Inches(0.2), Inches(0.3),
             f'{phase}: {weeks}', size=Pt(10), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(0.45), w - Inches(0.2), Inches(0.6),
             title, size=Pt(14), color=NAVY, bold=True)
    add_text(s, x + Inches(0.1), y + Inches(1.1), w - Inches(0.2), Inches(3.0),
             '\n'.join(items), size=Pt(10), color=BODY_GRAY)

# Timeline bar
by = Inches(5.8)
add_rect(s, Inches(0.6), by, Inches(11.6), Inches(0.2), fill_color=LIGHT_GRAY)
widths = [Inches(2.6), Inches(2.6), Inches(5.1), Inches(3.9)]
pcols = [GREEN, ACCENT_BLUE, MID_BLUE, NAVY]
cx_pos = Inches(0.6)
for pw, pc in zip(widths, pcols):
    add_rect(s, cx_pos, by, pw, Inches(0.2), fill_color=pc)
    cx_pos += pw

for wk, lbl in [(0, 'Wk 0'), (4, 'Wk 4'), (8, 'Wk 8'), (12, 'Wk 12'), (16, 'Wk 16'), (18, 'Wk 18')]:
    xp = Inches(0.6) + Inches(wk / 18.0 * 11.6)
    add_text(s, xp - Inches(0.2), by + Inches(0.25), Inches(0.5), Inches(0.2),
             lbl, size=Pt(8), color=BODY_GRAY, alignment=PP_ALIGN.CENTER)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 19: CONTINGENCY
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Contingency: If Rapid Ship Isn\u2019t Feasible')
add_accent_bar(s)
add_subheader(s, 'If NEITHER path achieves ~2-week FT')

ay = Inches(2.0)
# Option A
add_rect(s, Inches(0.8), ay, Inches(5.5), Inches(3.5), line_color=ACCENT_BLUE)
add_rect(s, Inches(0.8), ay, Inches(5.5), Inches(0.45), fill_color=ACCENT_BLUE)
add_text(s, Inches(1.0), ay + Inches(0.07), Inches(5.1), Inches(0.35),
         'OPTION A: COLLAPSE TO 2 TIERS', size=Pt(12), color=WHITE, bold=True)
add_text(s, Inches(1.0), ay + Inches(0.6), Inches(5.1), Inches(2.5),
         'Kold Locker + Fineline. Retire Fast-Trak.\n\nCreate \u201cFineline Standard\u201d pricing tier for 2-ft / standard-option orders at 10\u201312% discount vs. full custom.\n\nPros: Clean, simple, industry standard.\nCons: Loses FT brand. No mid-tier speed.',
         size=Pt(12), color=BODY_GRAY)

# Option B
add_rect(s, Inches(6.8), ay, Inches(5.5), Inches(3.5), line_color=ORANGE)
add_rect(s, Inches(6.8), ay, Inches(5.5), Inches(0.45), fill_color=ORANGE)
add_text(s, Inches(7.0), ay + Inches(0.07), Inches(5.1), Inches(0.35),
         'OPTION B: KEEP 3 TIERS, RENAME FT', size=Pt(12), color=WHITE, bold=True)
add_text(s, Inches(7.0), ay + Inches(0.6), Inches(5.1), Inches(2.5),
         'Drop \u201cFast\u201d from the name. Rename to emphasize flexibility/options (e.g., FlexLine, SelectLine, OptiKold).\n\nReposition as \u201cmore than KL, simpler than Fineline.\u201d\n\nCons: Middle tier = internal convenience, not customer value.',
         size=Pt(12), color=BODY_GRAY)

add_rect(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8), fill_color=ORANGE_BG)
add_text(s, Inches(1.0), Inches(5.85), Inches(11), Inches(0.7),
         'The 3-tier structure only earns its complexity if each tier delivers a distinct, tangible benefit.\nWithout speed, the middle tier is an internal convenience, not a customer value proposition.',
         size=Pt(13), color=NAVY, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 20: COLLATERAL ISSUES — FIX NOW
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Fix Now: Collateral Inconsistencies')
add_accent_bar(s)
add_subheader(s, 'These are not questions \u2014 they are immediate action items')

issues = [
    ('Door Widths', 'Comparison sheet: 30\u201d only for KL. Build Your KL flyer + website: 26\u201d/30\u201d/36\u201d.', 'Which is correct? Reconcile across all materials.'),
    ('Hinge Options', 'Comparison sheet: left only. Other materials: left/right available.', 'Confirm stocked set. Update comparison sheet.'),
    ('Refrigeration Naming', 'Website: \u201cSplit-Pak A2L.\u201d Comparison sheet: \u201cSplit-Pak Remote.\u201d', 'Standardize naming across all touchpoints.'),
    ('Feature Matrix', 'Walk-Ins brochure checkmarks make FT and Fineline look identical.', 'Rebuild matrix to show actual differentiation.'),
]
for i, (title, problem, action) in enumerate(issues):
    y = Inches(2.0) + Inches(i * 1.15)
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(1.0), fill_color=RED_BG if i % 2 == 0 else ORANGE_BG)
    add_rect(s, Inches(0.8), y + Inches(0.02), Inches(0.35), Inches(0.35), fill_color=RED)
    add_text(s, Inches(0.82), y - Inches(0.01), Inches(0.35), Inches(0.35),
             str(i + 1), font_name=FONT_DISPLAY, size=Pt(18), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(1.3), y + Inches(0.05), Inches(2.5), Inches(0.3),
             title, size=Pt(14), color=NAVY, bold=True)
    add_text(s, Inches(1.3), y + Inches(0.35), Inches(5.5), Inches(0.5),
             problem, size=Pt(11), color=BODY_GRAY)
    add_text(s, Inches(7.5), y + Inches(0.05), Inches(4.5), Inches(0.3),
             'ACTION:', size=Pt(9), color=RED, bold=True)
    add_text(s, Inches(7.5), y + Inches(0.3), Inches(4.5), Inches(0.6),
             action, size=Pt(11), color=NAVY, bold=True)

add_rect(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4), fill_color=NAVY)
add_text(s, Inches(1.0), Inches(6.32), Inches(11), Inches(0.35),
         'These undermine credibility regardless of which strategic direction we choose. Fix them first.',
         size=Pt(12), color=WHITE, bold=True)

add_logo_dark(s)
add_footer(s)

# ================================================================
# SLIDE 21: QUESTION SUMMARY / SCORECARD
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
add_header(s, 'Question Scorecard')
add_accent_bar(s)
add_subheader(s, 'All 43 questions at a glance \u2014 track answers here')

# Column headers
hdr_y = Inches(1.8)
add_text(s, Inches(0.5), hdr_y, Inches(0.3), Inches(0.25), '#', size=Pt(8), color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(0.85), hdr_y, Inches(6.2), Inches(0.25), 'QUESTION', size=Pt(8), color=NAVY, bold=True)
add_text(s, Inches(7.2), hdr_y, Inches(1.3), Inches(0.25), 'OWNER', size=Pt(8), color=NAVY, bold=True)
add_text(s, Inches(8.6), hdr_y, Inches(1.0), Inches(0.25), 'TYPE', size=Pt(8), color=NAVY, bold=True)
add_text(s, Inches(9.8), hdr_y, Inches(1.0), Inches(0.25), 'STATUS', size=Pt(8), color=NAVY, bold=True)
add_rect(s, Inches(0.5), hdr_y + Inches(0.25), Inches(10.5), Pt(1), fill_color=NAVY)

all_questions = [
    ('Q1', 'FT orders buildable as Fineline?', 'Sales/Ops', 'Data'),
    ('Q2', 'Dealer confusion FT vs Fineline?', 'Sales', 'Input'),
    ('Q3', 'Lost deals due to FT lead time?', 'Sales', 'Input'),
    ('Q4', 'Volume split KL/FT/Fineline?', 'Finance', 'Data'),
    ('Q5', 'Assembly in 10 biz days from stock?', 'Mfg', 'Y/N'),
    ('Q6', 'Filler panels for odd-ft sizes?', 'Eng', 'Y/N'),
    ('Q7', 'Plant capacity for filler panels?', 'Mfg', 'Y/N'),
    ('Q8', 'Assembly time delta KL vs FT?', 'Mfg', 'Data'),
    ('Q9', 'Why was rapid-build abandoned?', 'Ops', 'Input'),
    ('Q10', '% FT orders in 2-ft increments?', 'Ops', 'Data'),
    ('Q11', '% FT orders using odd-ft sizes?', 'Ops', 'Data'),
    ('Q12', 'Odd-ft price premium?', 'Sales', 'Data'),
    ('Q13', 'Additional component inventory?', 'Supply', 'Data'),
    ('Q14', 'Carrying cost of inventory?', 'Finance', 'Data'),
    ('Q15', 'Always-Fineline options?', 'Product', 'Input'),
    ('Q16', 'Non-rapid-shippable FT options?', 'Mfg', 'Input'),
    ('Q17', 'Fineline options movable to rapid?', 'Mfg', 'Input'),
    ('Q18', 'Right FT triggers? Others?', 'Product', 'Input'),
    ('Q19', '8\'7" in KL or FT?', 'Sales', 'Decision'),
    ('Q20', '% KL\u2192FT from sizing only?', 'Ops', 'Data'),
    ('Q21', 'Replacement price sensitivity?', 'Sales', 'Input'),
    ('Q22', 'Ecommerce price elasticity?', 'Ecomm', 'Data'),
    ('Q23', 'FT speed premium over Fineline?', 'Strategic', 'Decision'),
    ('Q24', 'Competitor pricing gaps?', 'Comp Intel', 'Data'),
    ('Q25', 'WebstaurantStore appetite?', 'Key Acct', 'Input'),
    ('Q26', 'FT ecommerce configs + premium?', 'Ecomm', 'Data'),
    ('Q27', 'Ship date = higher conversion?', 'Ecomm', 'Input'),
    ('Q28', 'Traditional dealer FT/FL view?', 'Sales', 'Input'),
    ('Q29', 'FT name equity with dealers?', 'Sales', 'Input'),
    ('Q30', 'Rename transition cost?', 'Marketing', 'Data'),
    ('Q31', 'Name change = momentum or confusion?', 'Sales', 'Input'),
    ('Q32', 'Master-Bilt Ready-Bilt lead time?', 'Comp Intel', 'Data'),
    ('Q33', 'Competitor tier positioning + pricing?', 'Comp Intel', 'Data'),
    ('Q34', 'Total valid FT configs / SKU count?', 'Eng/Product', 'Data'),
    ('Q35', 'Pre-engineer BOMs for all FT?', 'Engineering', 'Y/N'),
    ('Q36', 'Avg eng/quoting time per FT order?', 'Operations', 'Data'),
    ('Q37', 'Reduce to <500 SKUs w/o losing demand?', 'Product/Eng', 'Input'),
    ('Q38', 'Production lead time w/ zero eng?', 'Ops/Mfg', 'Data'),
    ('Q39', 'Process steps that add time?', 'Operations', 'Input'),
    ('Q40', 'Dedicated slots hit 2 wks reliably?', 'Ops/Mfg', 'Y/N'),
    ('Q41', 'Pre-stageable components + cost?', 'Supply/Fin', 'Data'),
    ('Q42', 'Path A vs B: lower total cost?', 'Finance/Ops', 'Analysis'),
    ('Q43', 'Hybrid feasible? A + B combo?', 'Strategic', 'Decision'),
]

row_h = Inches(0.112)
start_y = hdr_y + Inches(0.3)
for i, (qnum, qtext, owner, qtype) in enumerate(all_questions):
    y = start_y + row_h * i
    bg = LIGHT_BLUE_BG if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.5), y, Inches(10.5), row_h, fill_color=bg)
    add_text(s, Inches(0.5), y - Inches(0.01), Inches(0.3), row_h,
             qnum, size=Pt(6), color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(0.85), y - Inches(0.01), Inches(6.2), row_h,
             qtext, size=Pt(6), color=BODY_GRAY)
    add_text(s, Inches(7.2), y - Inches(0.01), Inches(1.3), row_h,
             owner, size=Pt(6), color=MID_BLUE, bold=True)
    add_text(s, Inches(8.6), y - Inches(0.01), Inches(1.0), row_h,
             qtype, size=Pt(6), color=BODY_GRAY)
    # Empty status box
    add_rect(s, Inches(9.8), y + Inches(0.01), Inches(0.7), row_h - Inches(0.02), line_color=MUTED_BLUE)

add_logo_dark(s)

# ================================================================
# SLIDE 22: CLOSING (dark)
# ================================================================
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill_color=NAVY)
s.shapes.add_picture(logo_white, Inches(1.0), Inches(1.5), Inches(4.0))

add_text(s, Inches(1.0), Inches(3.0), Inches(10), Inches(0.6),
         'Answer the Questions.',
         font_name=FONT_DISPLAY, size=Pt(44), color=WHITE, bold=True)
add_text(s, Inches(1.0), Inches(3.6), Inches(10), Inches(0.6),
         'Then Decide.',
         font_name=FONT_DISPLAY, size=Pt(44), color=ACCENT_BLUE, bold=True)

add_rect(s, Inches(1.0), Inches(4.3), Inches(2.0), Pt(2), fill_color=ACCENT_BLUE)

add_text(s, Inches(1.0), Inches(4.6), Inches(10), Inches(0.3),
         'IMMEDIATE ACTIONS', size=Pt(11), color=ACCENT_BLUE, bold=True)
add_multiline(s, Inches(1.0), Inches(4.9), Inches(10), Inches(1.2),
              ['\u2022 Pull FT order data (Q10\u2013Q12) \u2014 validates Path A sizing hypothesis',
               '\u2022 Schedule manufacturing feasibility assessment (Q5\u2013Q9) \u2014 Path A',
               '\u2022 Enumerate FT SKU count (Q34) + BOM effort (Q35) \u2014 Path B',
               '\u2022 Get Ops lead time with zero engineering (Q38) \u2014 Path B',
               '\u2022 Resolve 4 collateral inconsistencies \u2014 credibility cost today'],
              size=Pt(12), color=MUTED_BLUE)

add_text(s, Inches(1.0), Inches(6.3), Inches(10), Inches(0.5),
         '43 Questions  \u2192  Two Paths  \u2192  Data + Decisions  \u2192  Execute',
         font_name=FONT_DISPLAY, size=Pt(24), color=ACCENT_BLUE, bold=True)

add_text(s, Inches(1.0), Inches(6.8), Inches(5), Inches(0.3),
         'INTERNAL \u2014 CONFIDENTIAL  |  February 2026  |  v3', size=Pt(9), color=SUBTLE_BLUE)


# ================================================================
# SAVE
# ================================================================
output = r'C:\Users\andre\Desktop\Claude Code\Fast Trak\Fast-Trak-Strategy_v3.pptx'
prs.save(output)
print(f'Saved: {output}')
print(f'Slides: {len(prs.slides)}')
