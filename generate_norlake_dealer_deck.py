"""Generate Norlake Dealer Partner Deck — Benefits of In-House Refrigeration & Product Line Overview."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === BRAND CONFIGURATION (Norlake) ===
NAVY = RGBColor(0x00, 0x28, 0x57)
DARK_NAVY = RGBColor(0x00, 0x15, 0x32)
ACCENT_BLUE = RGBColor(0x2B, 0x7C, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF1, 0xED)
BODY_GRAY = RGBColor(0x2C, 0x3E, 0x50)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
MID_BLUE = RGBColor(0x1A, 0x4D, 0x7C)
ICE_BLUE = RGBColor(0xD6, 0xEE, 0xFF)
LIGHT_BLUE_BG = RGBColor(0xEB, 0xF5, 0xFF)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)

FONT_DISPLAY = 'Teko'
FONT_BODY = 'Trade Gothic Next'

ASSETS_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_WHITE = os.path.join(ASSETS_DIR, "assets", "logos", "norlake-no-oval-white.png")
LOGO_DARK = os.path.join(ASSETS_DIR, "assets", "logos", "norlake-no-oval-dark.png")

MARGIN_LEFT = Inches(0.75)
MARGIN_RIGHT = Inches(0.75)
CONTENT_WIDTH = Inches(11.833)

OUTPUT_PATH = r"C:\Users\andre\Desktop\Norlake-Dealer-Deck.pptx"


# === HELPER FUNCTIONS ===

def add_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, border=None):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    return shape


def add_text(slide, text, x, y, w=None, h=None, font=FONT_BODY, size=Pt(14),
             color=BODY_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a single-line text box."""
    w = w or Inches(11)
    h = h or Inches(0.5)
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def add_multiline(slide, lines, x, y, w, h=None, font=FONT_BODY, size=Pt(12),
                  color=BODY_GRAY, line_spacing=1.2, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a multi-paragraph text box from a list of lines."""
    h = h or Inches(3)
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(size.pt * (line_spacing - 1))
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
    return txBox


def add_header(slide, text, y=Inches(0.4)):
    """Add a BLUF headline at the top of a content slide."""
    add_text(slide, text, x=MARGIN_LEFT, y=y, w=Inches(10), h=Inches(0.7),
             font=FONT_DISPLAY, size=Pt(38), color=NAVY, bold=True)
    add_accent_bar(slide, x=MARGIN_LEFT, y=y + Inches(0.65))


def add_accent_bar(slide, x, y, width=Inches(1.5)):
    """Blue accent rule for visual rhythm."""
    add_rect(slide, x=x, y=y, w=width, h=Pt(3), fill=ACCENT_BLUE)


def add_footer(slide, text="CONFIDENTIAL"):
    """Standard footer bar with text and logo."""
    add_rect(slide, x=Inches(0), y=Inches(6.8), w=Inches(13.333), h=Inches(0.7), fill=NAVY)
    add_text(slide, text, x=MARGIN_LEFT, y=Inches(6.95), w=Inches(5), h=Inches(0.3),
             size=Pt(8), color=WHITE)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(10.5), Inches(6.85), width=Inches(2))


def add_stat_block(slide, number, label, x, y, number_size=Pt(56)):
    """Large number + label underneath."""
    add_text(slide, number, x=x, y=y, w=Inches(2.5), h=Inches(0.8),
             font=FONT_DISPLAY, size=number_size, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, label, x=x, y=y + Inches(0.75), w=Inches(2.5), h=Inches(0.5),
             font=FONT_BODY, size=Pt(11), color=BODY_GRAY,
             alignment=PP_ALIGN.CENTER)


def add_card(slide, title, body_lines, x, y, w, h, accent_color=ACCENT_BLUE):
    """Card with colored top border."""
    add_rect(slide, x=x, y=y, w=w, h=Inches(0.06), fill=accent_color)
    add_rect(slide, x=x, y=y + Inches(0.06), w=w, h=h - Inches(0.06), fill=WHITE, border=LIGHT_GRAY)
    add_text(slide, title, x=x + Inches(0.2), y=y + Inches(0.15), w=w - Inches(0.4), h=Inches(0.4),
             font=FONT_DISPLAY, size=Pt(22), color=NAVY, bold=True)
    add_multiline(slide, body_lines, x=x + Inches(0.2), y=y + Inches(0.55),
                  w=w - Inches(0.4), h=h - Inches(0.7), font=FONT_BODY, size=Pt(11), color=BODY_GRAY)


# === CREATE PRESENTATION ===
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# === SLIDE 1: TITLE — Norlake overview ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

# Title text
add_text(slide, "YOUR WALK-IN PARTNER", x=Inches(0.75), y=Inches(1.2),
         w=Inches(7), h=Inches(0.6), font=FONT_BODY, size=Pt(18),
         color=ACCENT_BLUE, bold=True)

add_text(slide, "14,000+ Configurations", x=Inches(0.75), y=Inches(1.8),
         w=Inches(7), h=Inches(0.9), font=FONT_DISPLAY, size=Pt(56),
         color=WHITE, bold=True)

add_text(slide, "Shipping in 2 Days from 660K Sq Ft", x=Inches(0.75), y=Inches(2.65),
         w=Inches(7), h=Inches(0.7), font=FONT_DISPLAY, size=Pt(40),
         color=WHITE, bold=True)

add_text(slide, "of US Manufacturing", x=Inches(0.75), y=Inches(3.25),
         w=Inches(7), h=Inches(0.7), font=FONT_DISPLAY, size=Pt(40),
         color=WHITE, bold=True)

# Subtitle
add_text(slide, "Dealer Partner Overview  |  In-House Refrigeration & Product Line",
         x=Inches(0.75), y=Inches(4.2), w=Inches(7), h=Inches(0.5),
         font=FONT_BODY, size=Pt(14), color=RGBColor(0x8A, 0xA8, 0xC8))

# Product image on right
img_path = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "kold-locker-capsule-pak-right.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(7.8), Inches(1.0), width=Inches(5))

# Logo at bottom (above safety line)
if os.path.exists(LOGO_WHITE):
    slide.shapes.add_picture(LOGO_WHITE, Inches(0.75), Inches(5.6), width=Inches(3))

# Date
add_text(slide, "April 2026", x=Inches(0.75), y=Inches(6.15), w=Inches(3), h=Inches(0.3),
         font=FONT_BODY, size=Pt(9), color=RGBColor(0x8A, 0xA8, 0xC8))


# === SLIDE 2: SECTION — Why In-House Refrigeration ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

add_text(slide, "WHY IN-HOUSE REFRIGERATION", x=Inches(0.75), y=Inches(2.5),
         w=Inches(11), h=Inches(0.9), font=FONT_DISPLAY, size=Pt(56),
         color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, x=Inches(5.7), y=Inches(3.35), width=Inches(2))

add_text(slide, "Simpler installs, lower costs, better efficiency",
         x=Inches(0.75), y=Inches(3.6), w=Inches(11), h=Inches(0.5),
         font=FONT_BODY, size=Pt(16), color=RGBColor(0x8A, 0xA8, 0xC8),
         alignment=PP_ALIGN.CENTER)

if os.path.exists(LOGO_WHITE):
    slide.shapes.add_picture(LOGO_WHITE, Inches(5.4), Inches(5.5), width=Inches(2.5))


# === SLIDE 3: Self-contained walk-ins eliminate complexity ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_header(slide, "Self-contained walk-ins eliminate mechanical rooms and cut install time")

# Three benefit cards
card_w = Inches(3.6)
card_h = Inches(3.5)
card_y = Inches(1.5)
gap = Inches(0.3)
start_x = Inches(0.75)

add_card(slide, "No Mechanical Room",
         ["Capsule Pak ECO mounts directly on the walk-in",
          "Frees up valuable floor space for revenue-generating use",
          "No remote condensing unit placement needed"],
         x=start_x, y=card_y, w=card_w, h=card_h, accent_color=ACCENT_BLUE)

add_card(slide, "Simplified Installation",
         ["Cord-and-plug electrical (most models 115V)",
          "Pre-charged, one-piece design \u2014 no field brazing",
          "Reduces install time and contractor coordination"],
         x=start_x + card_w + gap, y=card_y, w=card_w, h=card_h, accent_color=GREEN)

add_card(slide, "Lower Maintenance Costs",
         ["Single-point service \u2014 all components accessible on one unit",
          "No refrigerant line sets to leak or maintain",
          "LogiTemp controller enables remote diagnostics"],
         x=start_x + 2 * (card_w + gap), y=card_y, w=card_w, h=card_h, accent_color=ORANGE)

# Subtext
add_text(slide, "Capsule Pak ECO: patented self-contained system (U.S. Patent No. 11,859,885)",
         x=MARGIN_LEFT, y=Inches(5.2), w=Inches(10), h=Inches(0.4),
         font=FONT_BODY, size=Pt(10), color=MID_BLUE)

add_footer(slide)


# === SLIDE 4: R-290 Capsule Pak ECO energy savings ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_header(slide, "Capsule Pak ECO cuts energy costs ~50% with R-290 refrigerant")

# Left side: key benefits
bullets = [
    "R-290 natural refrigerant: GWP of 3 vs. R-404A GWP of 3,922",
    "Cooler energy: 2.5 kWh/day vs. 5.6 kWh/day legacy (56% reduction)",
    "Freezer energy: 5.6 kWh/day vs. 11.3 kWh/day legacy (51% reduction)",
    "Up to 28% shorter and 20% narrower than legacy units",
    "RSG transitioned to R-290 nearly 3 years ahead of AIM Act deadline",
]
add_multiline(slide, bullets, x=MARGIN_LEFT, y=Inches(1.5), w=Inches(5.8), h=Inches(3.0),
              font=FONT_BODY, size=Pt(13), color=BODY_GRAY, line_spacing=1.8)

# Callout box
add_rect(slide, x=Inches(0.75), y=Inches(4.8), w=Inches(5.8), h=Inches(1.0), fill=GREEN_BG, border=GREEN)
add_text(slide, "YOUR CUSTOMERS SAVE: ~50% on energy costs from day one",
         x=Inches(1.0), y=Inches(4.85), w=Inches(5.3), h=Inches(0.4),
         font=FONT_BODY, size=Pt(13), color=GREEN, bold=True)
add_text(slide, "Verified by independent third-party energy study",
         x=Inches(1.0), y=Inches(5.25), w=Inches(5.3), h=Inches(0.35),
         font=FONT_BODY, size=Pt(10), color=BODY_GRAY)

# Right side: product image
img_path = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "refrigeration", "capsule-pak-eco-transparent.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.3), width=Inches(5))

add_footer(slide)


# === SLIDE 5: STATS — Key numbers ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_header(slide, "Numbers your customers need to hear")

# Four stat blocks evenly spaced
stat_y = Inches(2.0)
stats = [
    ("14,000+", "Walk-in configurations\nin stock"),
    ("2-DAY", "Shipping on\nstock orders"),
    ("~50%", "Energy savings\nvs. legacy systems"),
    ("15-YEAR", "Structural panel\nwarranty"),
]

stat_w = Inches(2.8)
start_x = Inches(0.75)
gap = Inches(0.25)

for i, (number, label) in enumerate(stats):
    x = start_x + i * (stat_w + gap)
    # Background card
    add_rect(slide, x=x, y=stat_y, w=stat_w, h=Inches(3.0), fill=LIGHT_BLUE_BG)
    # Large number
    add_text(slide, number, x=x, y=stat_y + Inches(0.5), w=stat_w, h=Inches(0.9),
             font=FONT_DISPLAY, size=Pt(56), color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)
    # Label
    label_lines = label.split('\n')
    for li, line in enumerate(label_lines):
        add_text(slide, line, x=x, y=stat_y + Inches(1.5) + Inches(li * 0.3),
                 w=stat_w, h=Inches(0.3),
                 font=FONT_BODY, size=Pt(12), color=BODY_GRAY,
                 alignment=PP_ALIGN.CENTER)

add_footer(slide)


# === SLIDE 6: SECTION — The Norlake Product Line ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

add_text(slide, "THE NORLAKE PRODUCT LINE", x=Inches(0.75), y=Inches(2.5),
         w=Inches(11), h=Inches(0.9), font=FONT_DISPLAY, size=Pt(56),
         color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, x=Inches(5.7), y=Inches(3.35), width=Inches(2))

add_text(slide, "Walk-ins, refrigeration systems, and smart controls for every application",
         x=Inches(0.75), y=Inches(3.6), w=Inches(11), h=Inches(0.5),
         font=FONT_BODY, size=Pt(16), color=RGBColor(0x8A, 0xA8, 0xC8),
         alignment=PP_ALIGN.CENTER)

if os.path.exists(LOGO_WHITE):
    slide.shapes.add_picture(LOGO_WHITE, Inches(5.4), Inches(5.5), width=Inches(2.5))


# === SLIDE 7: Kold Locker — broadest walk-in range ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_header(slide, "Kold Locker: 14,000+ configs, 2-day ship, 15-year panel warranty")

# Left content
bullets = [
    "14,000+ configurations cover virtually any standard walk-in need",
    "Ships in 2 business days from stock inventory",
    "Foamed-in-place polyurethane insulation (R-32+ per inch)",
    "Cam-lock assembly for fast, reliable field installation",
    "18-month parts and labor warranty (50% above industry standard)",
]
add_multiline(slide, bullets, x=MARGIN_LEFT, y=Inches(1.5), w=Inches(6), h=Inches(3.0),
              font=FONT_BODY, size=Pt(13), color=BODY_GRAY, line_spacing=1.8)

# Key applications
add_text(slide, "KEY APPLICATIONS", x=MARGIN_LEFT, y=Inches(4.6), w=Inches(5), h=Inches(0.3),
         font=FONT_BODY, size=Pt(10), color=ACCENT_BLUE, bold=True)
add_text(slide, "Restaurants  |  Grocery  |  Institutional Kitchens  |  Beer Caves  |  Floral  |  Meat Processing",
         x=MARGIN_LEFT, y=Inches(4.9), w=Inches(6), h=Inches(0.4),
         font=FONT_BODY, size=Pt(11), color=BODY_GRAY)

# Right: product image (constrained to stay above safety line)
img_path = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "kold-locker-capsule-pak-left.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.5), height=Inches(4.5))

add_footer(slide)


# === SLIDE 8: Fast-Trak and FineLine — stock to custom to architectural ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_header(slide, "Fast-Trak and FineLine cover every project from stock to architectural")

# Two cards side by side
card_w = Inches(5.5)
card_h = Inches(2.0)
card_y = Inches(1.5)

# Fast-Trak card
add_rect(slide, x=Inches(0.75), y=card_y, w=card_w, h=Inches(0.06), fill=ACCENT_BLUE)
add_rect(slide, x=Inches(0.75), y=card_y + Inches(0.06), w=card_w, h=card_h - Inches(0.06),
         fill=WHITE, border=LIGHT_GRAY)
add_text(slide, "FAST-TRAK", x=Inches(0.95), y=card_y + Inches(0.15), w=Inches(5), h=Inches(0.4),
         font=FONT_DISPLAY, size=Pt(26), color=NAVY, bold=True)
ft_bullets = [
    "Custom and semi-custom walk-ins",
    "Non-standard dimensions and specialized layouts",
    "Configurations beyond Kold Locker catalog",
    "Same foamed-in-place panel quality",
]
add_multiline(slide, ft_bullets, x=Inches(0.95), y=card_y + Inches(0.6), w=Inches(4.8), h=Inches(2.5),
              font=FONT_BODY, size=Pt(12), color=BODY_GRAY, line_spacing=1.6)

# FineLine card
fl_x = Inches(6.75)
add_rect(slide, x=fl_x, y=card_y, w=card_w, h=Inches(0.06), fill=GOLD)
add_rect(slide, x=fl_x, y=card_y + Inches(0.06), w=card_w, h=card_h - Inches(0.06),
         fill=WHITE, border=LIGHT_GRAY)
add_text(slide, "FINELINE", x=fl_x + Inches(0.2), y=card_y + Inches(0.15), w=Inches(5), h=Inches(0.4),
         font=FONT_DISPLAY, size=Pt(26), color=NAVY, bold=True)
fl_bullets = [
    "Premium architectural walk-in line",
    "Design integration with visible spaces",
    "Aesthetics alongside full performance",
    "Ideal for front-of-house and retail environments",
]
add_multiline(slide, fl_bullets, x=fl_x + Inches(0.2), y=card_y + Inches(0.6), w=Inches(4.8), h=Inches(2.5),
              font=FONT_BODY, size=Pt(12), color=BODY_GRAY, line_spacing=1.6)

# Product images inside card area (constrained height to stay above safety line)
ft_img = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "fast-trak.jpg")
if os.path.exists(ft_img):
    slide.shapes.add_picture(ft_img, Inches(1.2), Inches(3.5), height=Inches(2.5))

fl_img = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "fineline.jpg")
if os.path.exists(fl_img):
    slide.shapes.add_picture(fl_img, Inches(7.2), Inches(3.5), height=Inches(2.5))

add_footer(slide)


# === SLIDE 9: LogiTemp smart controllers ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_header(slide, "LogiTemp controllers save up to 27% energy with smart defrost")

# Left content
bullets = [
    "Up to 27% energy savings vs. all-mechanical systems",
    "Demand Defrost: defrosts only as needed, not on a timer",
    "30-day room and coil temperature data logging",
    "Wireless or Cat5 remote monitoring and programming",
    "Standard on all Capsule Pak ECO and Split-Pak freezer 6 HP+ systems",
]
add_multiline(slide, bullets, x=MARGIN_LEFT, y=Inches(1.5), w=Inches(6), h=Inches(3.0),
              font=FONT_BODY, size=Pt(13), color=BODY_GRAY, line_spacing=1.8)

# Reverse Cycle Defrost callout
add_rect(slide, x=MARGIN_LEFT, y=Inches(4.7), w=Inches(6), h=Inches(1.2), fill=LIGHT_BLUE_BG, border=ACCENT_BLUE)
add_text(slide, "REVERSE CYCLE DEFROST OPTION (SPLIT-PAK 6 HP+)",
         x=Inches(1.0), y=Inches(4.75), w=Inches(5.5), h=Inches(0.3),
         font=FONT_BODY, size=Pt(10), color=ACCENT_BLUE, bold=True)
add_text(slide, "80% less defrost energy  |  Freezer defrost in 3-5 min vs. 20-30 min  |  93% faster cooler defrost",
         x=Inches(1.0), y=Inches(5.1), w=Inches(5.5), h=Inches(0.4),
         font=FONT_BODY, size=Pt(11), color=BODY_GRAY)

# Right: controller image
img_path = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "controllers", "logitemp-laptop.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.3), width=Inches(5))

add_footer(slide)


# === SLIDE 10: Split-Pak remote systems ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_header(slide, "Split-Pak remote systems extend coverage from 1/2 HP to 15 HP")

# Left content
bullets = [
    "Remote condensing units from 1/2 HP to 15 HP (contact factory for larger)",
    "Pre-charged quick-connect line sets up to 15 HP / 50 ft \u2014 reduces field brazing",
    "Pre-wired, factory-mounted components with traceable tubing",
    "EC motors standard on single-phase models",
    "All systems quoted for 2026+ use AIM Act-compliant refrigerants",
]
add_multiline(slide, bullets, x=MARGIN_LEFT, y=Inches(1.5), w=Inches(6), h=Inches(2.5),
              font=FONT_BODY, size=Pt(13), color=BODY_GRAY, line_spacing=1.8)

# Split-Pak ECO callout
add_rect(slide, x=MARGIN_LEFT, y=Inches(4.3), w=Inches(6), h=Inches(1.4), fill=GREEN_BG, border=GREEN)
add_text(slide, "SPLIT-PAK ECO: INDUSTRY'S FIRST R-290 REMOTE SYSTEM",
         x=Inches(1.0), y=Inches(4.35), w=Inches(5.5), h=Inches(0.3),
         font=FONT_BODY, size=Pt(10), color=GREEN, bold=True)
add_text(slide, "Extends natural refrigerant benefits beyond the 1.75 HP self-contained limit. Gives your customers R-290 efficiency for larger walk-in applications.",
         x=Inches(1.0), y=Inches(4.7), w=Inches(5.5), h=Inches(0.7),
         font=FONT_BODY, size=Pt(11), color=BODY_GRAY)

# Right: kitchen setting
img_path = os.path.join(ASSETS_DIR, "assets", "products", "norlake", "walk-ins", "kitchen-setting-left.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.3), width=Inches(5))

add_footer(slide)


# === SLIDE 11: STATS — Corporate credentials ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_header(slide, "85+ years, 660K sq ft, ISO certified, Made in USA")

# Four stat blocks
stat_y = Inches(1.8)
stats = [
    ("85+", "Combined years of\ncommercial refrigeration"),
    ("660K+", "Sq ft of manufacturing\nacross two US plants"),
    ("ISO 9001\n& 14001", "Quality + environmental\ncertification (DEKRA)"),
    ("MADE IN\nUSA", "Hudson, WI and\nNew Albany, MS"),
]

stat_w = Inches(2.8)
start_x = Inches(0.75)
gap = Inches(0.25)

for i, (number, label) in enumerate(stats):
    x = start_x + i * (stat_w + gap)
    add_rect(slide, x=x, y=stat_y, w=stat_w, h=Inches(2.8), fill=LIGHT_BLUE_BG)

    num_lines = number.split('\n')
    if len(num_lines) == 1:
        add_text(slide, number, x=x, y=stat_y + Inches(0.3), w=stat_w, h=Inches(0.9),
                 font=FONT_DISPLAY, size=Pt(56), color=ACCENT_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    else:
        for ni, nline in enumerate(num_lines):
            add_text(slide, nline, x=x, y=stat_y + Inches(0.15) + Inches(ni * 0.45),
                     w=stat_w, h=Inches(0.5),
                     font=FONT_DISPLAY, size=Pt(36), color=ACCENT_BLUE, bold=True,
                     alignment=PP_ALIGN.CENTER)

    label_lines = label.split('\n')
    for li, line in enumerate(label_lines):
        add_text(slide, line, x=x, y=stat_y + Inches(1.4) + Inches(li * 0.3),
                 w=stat_w, h=Inches(0.3),
                 font=FONT_BODY, size=Pt(12), color=BODY_GRAY,
                 alignment=PP_ALIGN.CENTER)

# ISO logos (above safety line)
iso_9001 = os.path.join(ASSETS_DIR, "assets", "logos", "iso-9001.png")
iso_14001 = os.path.join(ASSETS_DIR, "assets", "logos", "iso-14001.png")
if os.path.exists(iso_9001):
    slide.shapes.add_picture(iso_9001, Inches(4.8), Inches(4.9), width=Inches(1.0))
if os.path.exists(iso_14001):
    slide.shapes.add_picture(iso_14001, Inches(6.2), Inches(4.9), width=Inches(1.0))

add_footer(slide)


# === SLIDE 12: CLOSER — Partner with Norlake ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

# Main CTA
add_text(slide, "PARTNER WITH NORLAKE", x=Inches(0.75), y=Inches(1.0),
         w=Inches(11), h=Inches(0.9), font=FONT_DISPLAY, size=Pt(56),
         color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, x=Inches(5.7), y=Inches(1.9), width=Inches(2))

add_text(slide, "The most efficient, available, and well-supported walk-in refrigeration in the industry",
         x=Inches(1.5), y=Inches(2.2), w=Inches(10), h=Inches(0.6),
         font=FONT_BODY, size=Pt(18), color=RGBColor(0x8A, 0xA8, 0xC8),
         alignment=PP_ALIGN.CENTER)

# Three selling points
selling_points = [
    ("14,000+ CONFIGS", "2-day stock shipping"),
    ("~50% ENERGY SAVINGS", "R-290 Capsule Pak ECO"),
    ("15-YEAR WARRANTY", "Foamed-in-place panels"),
]

sp_y = Inches(3.3)
sp_w = Inches(3.2)
sp_gap = Inches(0.5)
sp_start = Inches(1.5)

for i, (title, subtitle) in enumerate(selling_points):
    x = sp_start + i * (sp_w + sp_gap)
    add_rect(slide, x=x, y=sp_y, w=sp_w, h=Inches(1.2), fill=RGBColor(0x0B, 0x35, 0x65))
    add_text(slide, title, x=x, y=sp_y + Inches(0.15), w=sp_w, h=Inches(0.4),
             font=FONT_DISPLAY, size=Pt(24), color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, subtitle, x=x, y=sp_y + Inches(0.6), w=sp_w, h=Inches(0.3),
             font=FONT_BODY, size=Pt(12), color=RGBColor(0x8A, 0xA8, 0xC8),
             alignment=PP_ALIGN.CENTER)

# Contact info
add_text(slide, "Contact your Norlake representative  |  800-955-5253  |  norlake.com",
         x=Inches(0.75), y=Inches(5.0), w=Inches(11), h=Inches(0.4),
         font=FONT_BODY, size=Pt(14), color=WHITE, alignment=PP_ALIGN.CENTER)

# Logo (above safety line)
if os.path.exists(LOGO_WHITE):
    slide.shapes.add_picture(LOGO_WHITE, Inches(5.2), Inches(5.6), width=Inches(3))


# === SAVE ===
prs.save(OUTPUT_PATH)
print(f"Saved: {OUTPUT_PATH}")
